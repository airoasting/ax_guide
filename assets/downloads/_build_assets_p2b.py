"""
Build P2-B (Strategy track, batch 2) download assets for AX Guide for Enterprise.

Outputs (in same directory as this script):
- incident-scenarios.docx        (strategy-3)
- job-redesign-workbook.docx     (strategy-4)
- operating-models.pptx          (strategy-4)
- capability-map.xlsx            (strategy-4)
- vendor-rfp-template.docx       (strategy-5)
- architecture-patterns.pptx     (strategy-5)
- tco-comparison.xlsx            (strategy-5)

Run from repo root:
    python3 assets/downloads/_build_assets_p2b.py
"""

import os

from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from pptx import Presentation
from pptx.util import Inches as Pin, Pt as PPt
from pptx.dml.color import RGBColor as PColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTDIR = SCRIPT_DIR
FOOTER = "AX Guide for Enterprise · v0.2 / jade@linercorp.com"


# ---------- common style helpers (mirror _build_assets.py) ----------

THIN = Side(border_style="thin", color="CCCCCC")
BORDER = Border(left=THIN, right=THIN, top=THIN, bottom=THIN)

HEADER_FILL = PatternFill("solid", fgColor="1F2937")
SUBHEADER_FILL = PatternFill("solid", fgColor="E5E7EB")
ACCENT_FILL = PatternFill("solid", fgColor="FEF3C7")
GOOD_FILL = PatternFill("solid", fgColor="D1FAE5")
WARN_FILL = PatternFill("solid", fgColor="FEE2E2")

HEADER_FONT = Font(name="Inter", size=12, bold=True, color="FFFFFF")
SUBHEADER_FONT = Font(name="Inter", size=11, bold=True, color="111827")
BODY_FONT = Font(name="Inter", size=11, color="111827")
HINT_FONT = Font(name="Inter", size=10, italic=True, color="6B7280")
FOOTER_FONT = Font(name="Inter", size=9, italic=True, color="6B7280")


def set_col_widths(ws, widths):
    for i, w in enumerate(widths, start=1):
        ws.column_dimensions[get_column_letter(i)].width = w


def put(ws, cell, value, font=None, fill=None, border=None, align=None):
    c = ws[cell]
    c.value = value
    if font:
        c.font = font
    if fill:
        c.fill = fill
    if border:
        c.border = border
    if align:
        c.alignment = align


def _add_page_number_footer(doc, footer_text):
    """Add page numbering + custom footer text to the document."""
    section = doc.sections[0]
    footer = section.footer
    p = footer.paragraphs[0]
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    run = p.add_run(f"{footer_text}    ·    ")
    run.font.size = Pt(9)
    run.font.italic = True
    run.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

    fld_run = p.add_run()
    fld_run.font.size = Pt(9)
    fld_run.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

    fldChar1 = OxmlElement("w:fldChar")
    fldChar1.set(qn("w:fldCharType"), "begin")
    instrText = OxmlElement("w:instrText")
    instrText.set(qn("xml:space"), "preserve")
    instrText.text = "PAGE"
    fldChar2 = OxmlElement("w:fldChar")
    fldChar2.set(qn("w:fldCharType"), "end")
    fld_run._r.append(fldChar1)
    fld_run._r.append(instrText)
    fld_run._r.append(fldChar2)

    fld_run.add_text(" / ")

    fldChar3 = OxmlElement("w:fldChar")
    fldChar3.set(qn("w:fldCharType"), "begin")
    instrText2 = OxmlElement("w:instrText")
    instrText2.set(qn("xml:space"), "preserve")
    instrText2.text = "NUMPAGES"
    fldChar4 = OxmlElement("w:fldChar")
    fldChar4.set(qn("w:fldCharType"), "end")
    fld_run._r.append(fldChar3)
    fld_run._r.append(instrText2)
    fld_run._r.append(fldChar4)


def _docx_set_margins(doc, inches=1.0):
    for section in doc.sections:
        section.top_margin = Inches(inches)
        section.bottom_margin = Inches(inches)
        section.left_margin = Inches(inches)
        section.right_margin = Inches(inches)


def _docx_title(doc, title_text, meta_text=None):
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run(title_text)
    run.bold = True
    run.font.size = Pt(22)
    run.font.color.rgb = RGBColor(0x11, 0x18, 0x27)

    if meta_text:
        meta = doc.add_paragraph()
        meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
        mr = meta.add_run(meta_text)
        mr.font.size = Pt(10)
        mr.font.italic = True
        mr.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
    doc.add_paragraph()


def _bold_header_row(row_cells):
    for cell in row_cells:
        for p in cell.paragraphs:
            for r in p.runs:
                r.bold = True


# ---------- 1) incident-scenarios.docx ----------

def make_incident_scenarios():
    doc = Document()
    _docx_set_margins(doc)
    _docx_title(
        doc,
        "사고 시나리오 워크북",
        "문서번호: [INC-AI-001]    ·    적용일: [YYYY-MM-DD]    ·    소관부서: [정보보호팀 / AX TF]",
    )

    # 1. 목적
    doc.add_heading("1. 목적", level=1)
    doc.add_paragraph(
        "AI 사용 중 발생할 수 있는 사고 시나리오를 사전에 시뮬레이션하여 대응 역량을 강화합니다. "
        "각 시나리오는 인지부터 재발 방지까지 5단계 대응 절차를 가지며, 분기 1회 시뮬레이션을 권장합니다."
    )

    # Common 5-stage template factory
    stages = ["인지", "영향 평가", "고객/내부 통보", "사후 분석", "재발 방지"]

    def add_scenario(letter, title, summary, plan_lines):
        doc.add_heading(f"2{letter}. 시나리오 {letter[0].upper()}: {title}", level=1)
        doc.add_paragraph(summary)

        tbl = doc.add_table(rows=6, cols=3)
        tbl.style = "Light Grid Accent 1"
        hdr = tbl.rows[0].cells
        hdr[0].text = "단계"
        hdr[1].text = "행동"
        hdr[2].text = "담당 / 시한"
        _bold_header_row(hdr)
        for i, (stage_name, action) in enumerate(zip(stages, plan_lines), start=1):
            row = tbl.rows[i].cells
            row[0].text = f"({i}) {stage_name}"
            row[1].text = action
            row[2].text = "[ 담당자 ] / [ 시한 ]"
        doc.add_paragraph()

    # 시나리오 A
    add_scenario(
        "1",
        "외부 송출 사고",
        "직원이 AI가 생성한 출력물을 검증 없이 고객에게 그대로 송출하여 사실 오류·법적 리스크가 발생한 경우.",
        [
            "수신 고객 또는 내부 모니터링을 통해 오류 사실을 즉시 보고받고 사고 티켓을 발행합니다.",
            "오류의 사실 관계, 영향 범위(고객 수·금액·법적 리스크), 외부 노출 정도를 30분 이내 1차 평가합니다.",
            "고객에게 사실 인정·정정·향후 조치를 명시한 공식 안내를 발송하고, 필요 시 임원 승인 후 공개 사과합니다.",
            "근본 원인을 분석합니다 — 검증 절차 부재, 도구 선택 오류, 직원 교육 미흡 중 무엇이 1차 원인인지 식별합니다.",
            "검증 체크리스트 의무화, 송출 전 2인 리뷰, 해당 도구 사용 가이드 갱신, 분기 시뮬레이션 항목에 추가합니다.",
        ],
    )

    # 시나리오 B
    add_scenario(
        "2",
        "개인정보 노출",
        "직원이 고객 데이터(이메일, 전화번호, 주민번호 등)를 외부 LLM에 입력하여 개인정보가 외부로 송출된 경우.",
        [
            "DLP 알람, 동료 신고, 또는 본인 자진 신고를 통해 사고를 즉시 인지하고 정보보호팀에 1시간 이내 보고합니다.",
            "유출된 정보 유형·건수, 사용된 도구, 데이터 학습 가능성을 평가합니다. 개인정보보호위원회 신고 요건을 확인합니다.",
            "법적 신고 의무가 발생하면 72시간 이내 감독기관 신고, 고객에게는 영향이 있는 경우에만 정확한 사실을 통보합니다.",
            "해당 도구의 데이터 처리 정책을 재검토하고, 직원이 왜 외부 도구를 사용했는지(승인 도구 부재? 편의성?) 원인을 분석합니다.",
            "사내 승인 도구로 동일 작업이 가능한 경로를 마련하고, DLP 정책을 강화하며 해당 직원·팀에 추가 교육을 실시합니다.",
        ],
    )

    # 시나리오 C
    add_scenario(
        "3",
        "잘못된 자동화",
        "자동화 워크플로우(에이전트·스크립트)가 오작동하여 잘못된 결정을 실행한 경우. 예: 환불 자동 승인, 잘못된 고객 분류 발송 등.",
        [
            "모니터링 알람 또는 고객 클레임으로 사고를 인지하고 즉시 해당 워크플로우를 중단합니다.",
            "잘못된 결정의 건수·금액·고객 영향 범위를 데이터로 산출합니다. 자동 롤백이 가능한지 확인합니다.",
            "영향 받은 고객에게 정정 안내 및 보상 절차를 통보하고, 필요 시 임원·이사회에 보고합니다.",
            "워크플로우의 로직·테스트 커버리지·휴먼 인 더 루프(HITL) 부재 여부를 점검합니다.",
            "임계값 초과 시 자동 중단 가드, 일일 결과 리뷰, 카나리 배포 도입, 분기 시뮬레이션 항목에 추가합니다.",
        ],
    )

    # 5. 시뮬레이션 가이드
    doc.add_heading("3. 시뮬레이션 가이드", level=1)
    doc.add_paragraph(
        "분기 1회 시뮬레이션을 실시합니다. 각 시나리오를 한 번씩 돌리되, 매 분기 1개 시나리오를 선정해 "
        "실제 인원·시간·의사결정 경로를 점검합니다. RACI 매트릭스로 역할을 명확히 합니다."
    )

    raci = doc.add_table(rows=6, cols=5)
    raci.style = "Light Grid Accent 1"
    hdr = raci.rows[0].cells
    for i, h in enumerate(["역할", "인지", "영향 평가", "통보", "재발 방지"]):
        hdr[i].text = h
    _bold_header_row(hdr)
    raci_rows = [
        ("현업 발견자", "R", "C", "I", "I"),
        ("정보보호팀", "A", "R", "C", "C"),
        ("법무팀", "I", "C", "A", "I"),
        ("AX TF", "I", "C", "I", "R/A"),
        ("임원", "I", "I", "A (외부 공개 시)", "A"),
    ]
    for i, row in enumerate(raci_rows, start=1):
        for j, v in enumerate(row):
            raci.rows[i].cells[j].text = v
    doc.add_paragraph()
    doc.add_paragraph(
        "R = Responsible (실행), A = Accountable (최종 책임), C = Consulted (협의), I = Informed (통보)"
    )

    # 6. 빈 칸 워크시트
    doc.add_heading("4. 우리 회사 시나리오 (직접 작성)", level=1)
    doc.add_paragraph(
        "위 3개 시나리오 외에 우리 회사·우리 산업에서 발생 가능성이 가장 높은 시나리오 1개를 추가로 작성합니다."
    )

    custom = doc.add_table(rows=6, cols=3)
    custom.style = "Light Grid Accent 1"
    hdr = custom.rows[0].cells
    hdr[0].text = "단계"
    hdr[1].text = "행동"
    hdr[2].text = "담당 / 시한"
    _bold_header_row(hdr)
    for i, stage_name in enumerate(stages, start=1):
        row = custom.rows[i].cells
        row[0].text = f"({i}) {stage_name}"
        row[1].text = ""
        row[2].text = ""

    doc.add_paragraph()
    doc.add_paragraph(
        "[ 시나리오 제목 ]: ____________________________________________________"
    )
    doc.add_paragraph(
        "[ 가정 상황 ]: ______________________________________________________"
    )

    _add_page_number_footer(doc, FOOTER)

    out = os.path.join(OUTDIR, "incident-scenarios.docx")
    doc.save(out)
    print(f"  wrote {out}")


# ---------- 2) job-redesign-workbook.docx ----------

def make_job_redesign_workbook():
    doc = Document()
    _docx_set_margins(doc)
    _docx_title(
        doc,
        "직무 재정의 워크북",
        "문서번호: [JOB-AI-001]    ·    적용일: [YYYY-MM-DD]    ·    소관부서: [HR / AX TF]",
    )

    # 1. 목적
    doc.add_heading("1. 목적", level=1)
    doc.add_paragraph(
        "AX 시대에 직무가 어떻게 변하는지 직군별로 정리합니다. "
        "임원 회의 + 부서장 워크숍 자료로 사용하며, 각 직군의 변화·축소·생성·역량 4가지 축으로 표준화합니다."
    )

    # 직군별 예시
    job_examples = [
        (
            "경영기획·전략",
            "보고서 1차 초안 작성, 단순 데이터 정리",  # 변하는 일 (예시)
            "수작업 시장 리서치 정리, 반복적 차트 가공",  # 줄어드는 일
            "AI 출력 검증, 시나리오 분석 설계, 의사결정 품질 관리",  # 새로 생기는 일
            "데이터 해석력, 가설 설계, AI 결과 비판적 검증",  # 필요한 역량
        ),
        (
            "영업",
            "리드 리서치, 제안서 1차 초안, 이메일 작성",
            "단순 CRM 입력, 반복 팔로업 메일 작성",
            "AI 산출 검증, 고객 관계 심화, 복잡 협상 설계",
            "고객 인사이트 해석, 도구 활용, 협상력",
        ),
        (
            "마케팅",
            "콘텐츠 초안, 카피 변형, SEO·GEO 가공",
            "단순 소재 변형, 일반적 키워드 리서치",
            "브랜드 전략, 데이터 분석 설계, AI 콘텐츠 품질 관리",
            "브랜드 감각, 데이터 분석, AI 도구 큐레이션",
        ),
        (
            "HR",
            "JD 작성, 면접 평가 통합, 사내 안내문 작성",
            "반복적 스크리닝, 단순 인사 데이터 정리",
            "조직 진단 설계, 인재 전략, AI 채용 공정성 관리",
            "조직 진단, 공정성 감각, AI 활용 가이드 설계",
        ),
        (
            "재무·회계",
            "공시 분석, 결산 보조 자료, 보고서 1차 초안",
            "반복적 분개 점검, 단순 비교 분석",
            "재무 시나리오 설계, AI 결과 검증, 감사 대응",
            "회계 원칙, 데이터 검증, 규제 이해",
        ),
        (
            "개발",
            "코드 1차 생성, 테스트 작성, 문서화",
            "단순 보일러플레이트, 반복 리뷰",
            "아키텍처 설계, 코드 품질 가드레일, 사내 코파일럿 운영",
            "시스템 설계, 코드 리뷰력, AI 코드 검증",
        ),
        (
            "디자인",
            "시안 변형, 마이크로카피, 무드보드 생성",
            "단순 리사이즈, 반복적 톤 변형",
            "디자인 시스템 + AI 통합, 프로토타이핑 자동화 설계",
            "브랜드 일관성, 시스템 사고, AI 도구 큐레이션",
        ),
        (
            "고객지원",
            "응답 1차 초안, FAQ 작성, 통화 요약",
            "반복 응답 작성, 단순 분류·라우팅",
            "보이스봇 시나리오 설계, 고객 인사이트 분석, 에스컬레이션 품질 관리",
            "공감력, 시나리오 설계, AI 응답 품질 검증",
        ),
    ]

    rows_def = ["변하는 일", "줄어드는 일", "새로 생기는 일", "필요한 역량"]

    for idx, (job_name, ex_change, ex_reduce, ex_new, ex_skill) in enumerate(job_examples, start=1):
        doc.add_heading(f"2.{idx} {job_name}", level=1)
        examples = [ex_change, ex_reduce, ex_new, ex_skill]
        tbl = doc.add_table(rows=5, cols=2)
        tbl.style = "Light Grid Accent 1"
        hdr = tbl.rows[0].cells
        hdr[0].text = "항목"
        hdr[1].text = "내용"
        _bold_header_row(hdr)
        for i, (label, ex) in enumerate(zip(rows_def, examples), start=1):
            row = tbl.rows[i].cells
            row[0].text = label
            # 첫 줄은 예시 1개만 채우고 나머지는 빈 칸으로 → "예: ..." 표시
            row[1].text = f"예: {ex}\n· \n· "
        doc.add_paragraph()

    # 3. 워크숍 진행 가이드
    doc.add_heading("3. 워크숍 진행 가이드 (2시간)", level=1)
    guide = doc.add_table(rows=4, cols=3)
    guide.style = "Light Grid Accent 1"
    hdr = guide.rows[0].cells
    for i, h in enumerate(["시간", "내용", "산출물"]):
        hdr[i].text = h
    _bold_header_row(hdr)
    schedule = [
        ("0:00 - 0:30", "도입: AX 변화의 큰 그림 공유, 워크북 사용법 설명, 8개 직군 매핑 확인",
         "공유된 컨텍스트"),
        ("0:30 - 1:30", "직군별 그룹 작업: 각 직군 3-5명이 한 표를 함께 채움 (변하는 일·줄어드는 일·새로 생기는 일·필요한 역량)",
         "직군별 4축 표 1개씩"),
        ("1:30 - 2:00", "공유 및 정렬: 각 그룹 5분 발표, 공통 패턴·갈등 지점 정리",
         "통합 직무 변화 맵 1장"),
    ]
    for i, row in enumerate(schedule, start=1):
        for j, v in enumerate(row):
            guide.rows[i].cells[j].text = v

    doc.add_paragraph()
    doc.add_paragraph(
        "퍼실리테이터 팁: 첫 번째 칸은 예시를 보여주되, 나머지 칸은 반드시 참가자가 자신의 단어로 채우게 합니다. "
        "예시를 그대로 베끼면 워크숍의 의미가 사라집니다."
    )

    _add_page_number_footer(doc, FOOTER)

    out = os.path.join(OUTDIR, "job-redesign-workbook.docx")
    doc.save(out)
    print(f"  wrote {out}")


# ---------- 3) operating-models.pptx ----------

def make_operating_models():
    prs = Presentation()
    prs.slide_width = Pin(13.333)
    prs.slide_height = Pin(7.5)

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Title
    title_box = slide.shapes.add_textbox(Pin(0.5), Pin(0.3), Pin(12.3), Pin(0.7))
    p = title_box.text_frame.paragraphs[0]
    p.text = "5가지 AX 운영 모델 비교"
    p.font.size = PPt(26)
    p.font.bold = True
    p.font.color.rgb = PColor(0x11, 0x18, 0x27)
    p.alignment = PP_ALIGN.LEFT

    sub_box = slide.shapes.add_textbox(Pin(0.5), Pin(1.0), Pin(12.3), Pin(0.4))
    sp = sub_box.text_frame.paragraphs[0]
    sp.text = "회사 규모·문화·리스크 허용도에 따라 운영 모델을 선택합니다."
    sp.font.size = PPt(12)
    sp.font.color.rgb = PColor(0x6B, 0x72, 0x80)

    # Table
    headers = ["항목", "중앙 CoE", "임베디드 챔피언", "길드", "분산 자율", "하이브리드"]
    data = [
        ["구조", "별도 부서", "부서마다 1명", "공식 모임", "자유", "혼합"],
        ["의사결정", "중앙", "부서장", "합의", "개인", "중앙+부서"],
        ["적합 규모", "1000명+", "200-1000", "모든 규모", "<200", "500-2000"],
        ["장점", "일관성", "현장성", "학습 공유", "속도", "균형"],
        ["단점", "사일로", "표준화 약", "책임 분산", "사고 위험", "복잡"],
    ]

    rows = len(data) + 1
    cols = len(headers)
    left = Pin(0.5)
    top = Pin(1.6)
    width = Pin(12.3)
    height = Pin(4.8)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # column widths — first column slightly wider
    first_col = 2.0
    other = (12.3 - first_col) / 5
    table.columns[0].width = Pin(first_col)
    for c in range(1, cols):
        table.columns[c].width = Pin(other)

    # Header row
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Pin(0.1)
        tf.margin_right = Pin(0.1)
        para = tf.paragraphs[0]
        para.text = h
        para.alignment = PP_ALIGN.CENTER
        para.font.size = PPt(13)
        para.font.bold = True
        para.font.color.rgb = PColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PColor(0x1F, 0x29, 0x37)

    # Data rows
    for r, row_data in enumerate(data, start=1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Pin(0.1)
            tf.margin_right = Pin(0.1)
            para = tf.paragraphs[0]
            para.text = val
            para.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            para.font.size = PPt(11)
            para.font.bold = (c == 0)
            para.font.color.rgb = PColor(0x11, 0x18, 0x27)
            # alternate row shading
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PColor(0xF9, 0xFA, 0xFB)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PColor(0xFF, 0xFF, 0xFF)

    # Note
    note_box = slide.shapes.add_textbox(Pin(0.5), Pin(6.55), Pin(12.3), Pin(0.4))
    np_ = note_box.text_frame.paragraphs[0]
    np_.text = "단일 모델을 고집하지 마세요. 1년차는 임베디드 챔피언으로 시작 → 2년차에 하이브리드로 진화하는 경로가 일반적입니다."
    np_.font.size = PPt(10)
    np_.font.italic = True
    np_.font.color.rgb = PColor(0x6B, 0x72, 0x80)

    # Footer
    footer_box = slide.shapes.add_textbox(Pin(0.5), Pin(7.1), Pin(12.3), Pin(0.3))
    fp = footer_box.text_frame.paragraphs[0]
    fp.text = FOOTER
    fp.font.size = PPt(9)
    fp.font.italic = True
    fp.font.color.rgb = PColor(0x6B, 0x72, 0x80)
    fp.alignment = PP_ALIGN.RIGHT

    out = os.path.join(OUTDIR, "operating-models.pptx")
    prs.save(out)
    print(f"  wrote {out}")


# ---------- 4) capability-map.xlsx ----------

def make_capability_map():
    wb = Workbook()

    # Sheet 1: AX 역량 맵
    ws = wb.active
    ws.title = "AX 역량 맵"
    set_col_widths(ws, [18, 28, 28, 28, 28])

    put(ws, "A1", "AX 역량 맵 (직군 × 레벨)",
        Font(name="Inter", size=16, bold=True, color="111827"))
    ws.merge_cells("A1:E1")
    ws.row_dimensions[1].height = 26

    put(ws, "A2",
        "각 직군의 Lv1 ~ Lv4 정의입니다. 사내 표준으로 채택 후 부서별 평가에 사용하세요.",
        HINT_FONT)
    ws.merge_cells("A2:E2")

    headers = ["직군", "Lv1 Aware", "Lv2 User", "Lv3 Designer", "Lv4 Architect"]
    center = Alignment(horizontal="center", vertical="center", wrap_text=True)
    left = Alignment(horizontal="left", vertical="center", wrap_text=True)

    for i, h in enumerate(headers):
        put(ws, f"{get_column_letter(i+1)}4", h, HEADER_FONT, HEADER_FILL, BORDER, center)
    ws.row_dimensions[4].height = 24

    rows = [
        ("경영기획·전략",
         "AI 회의 참여 가능",
         "보고서 1편 AI로 작성",
         "시나리오 분석 자동화 설계",
         "전사 AX 비전 수립"),
        ("영업",
         "도구 1개 시연 가능",
         "리드 리서치 자동화",
         "영업 워크플로우 재설계",
         "CRM-AI 통합 아키텍처"),
        ("마케팅",
         "콘텐츠 1편 AI 가공",
         "주간 콘텐츠 50% AI 협업",
         "캠페인 분석 파이프라인",
         "GEO 전략 + 자동화"),
        ("HR",
         "JD 작성 AI 활용",
         "면접 평가 통합",
         "조직 진단 인터뷰 자동화",
         "사내 LLM RAG 운영"),
        ("재무·회계",
         "결산 보조 사용",
         "공시 분석 정형화",
         "예산 분석 자동화",
         "재무 데이터 RAG"),
        ("개발",
         "코드 리뷰 보조",
         "테스트 50% 자동 생성",
         "마이그레이션 자동화",
         "사내 코파일럿 운영"),
        ("디자인",
         "시안 변형 활용",
         "마이크로카피 자동",
         "프로토타이핑 자동화",
         "디자인 시스템 + AI"),
        ("고객지원",
         "응대 초안 활용",
         "FAQ 자동 갱신",
         "보이스봇 시나리오 설계",
         "고객지원 LLM 운영"),
    ]

    start = 5
    for i, row in enumerate(rows):
        r = start + i
        put(ws, f"A{r}", row[0], SUBHEADER_FONT, SUBHEADER_FILL, BORDER, center)
        for j, val in enumerate(row[1:], start=2):
            col = get_column_letter(j)
            put(ws, f"{col}{r}", val, BODY_FONT, None, BORDER, left)
        ws.row_dimensions[r].height = 42

    note_row = start + len(rows) + 2
    put(ws, f"A{note_row}",
        "Lv1: 인식 / Lv2: 사용자 / Lv3: 설계자 / Lv4: 아키텍트. 한 사람이 모든 직군 Lv4를 달성할 필요는 없습니다.",
        HINT_FONT, align=Alignment(wrap_text=True, vertical="top"))
    ws.merge_cells(f"A{note_row}:E{note_row}")

    footer_row = note_row + 2
    put(ws, f"A{footer_row}", FOOTER, FOOTER_FONT)
    ws.merge_cells(f"A{footer_row}:E{footer_row}")

    # Sheet 2: 현재 수준 평가
    ws2 = wb.create_sheet("현재 수준 평가")
    set_col_widths(ws2, [18, 22, 22, 22, 22, 28])

    put(ws2, "A1", "현재 수준 평가 (워크숍용)",
        Font(name="Inter", size=16, bold=True, color="111827"))
    ws2.merge_cells("A1:F1")
    ws2.row_dimensions[1].height = 26

    put(ws2, "A2",
        "각 직군의 Lv1~Lv4 도달 인원수 또는 평가자 이름·점수를 입력하세요. 예: '김OO 3점' 또는 '12명'.",
        HINT_FONT)
    ws2.merge_cells("A2:F2")

    headers2 = ["직군", "Lv1 Aware", "Lv2 User", "Lv3 Designer", "Lv4 Architect", "메모"]
    for i, h in enumerate(headers2):
        put(ws2, f"{get_column_letter(i+1)}4", h, HEADER_FONT, HEADER_FILL, BORDER, center)
    ws2.row_dimensions[4].height = 24

    jobs = [r[0] for r in rows]
    for i, job in enumerate(jobs):
        r = 5 + i
        put(ws2, f"A{r}", job, SUBHEADER_FONT, SUBHEADER_FILL, BORDER, center)
        for c in range(2, 6):
            col = get_column_letter(c)
            put(ws2, f"{col}{r}", "", BODY_FONT, None, BORDER, center)
        put(ws2, f"F{r}", "", BODY_FONT, None, BORDER, left)
        ws2.row_dimensions[r].height = 32

    note_row2 = 5 + len(jobs) + 2
    put(ws2, f"A{note_row2}",
        "워크숍 권장: 직군별 부서장이 자기 부서를 평가 → AX TF가 검증 → 분기마다 갱신.",
        HINT_FONT)
    ws2.merge_cells(f"A{note_row2}:F{note_row2}")

    put(ws2, f"A{note_row2+2}", FOOTER, FOOTER_FONT)
    ws2.merge_cells(f"A{note_row2+2}:F{note_row2+2}")

    out = os.path.join(OUTDIR, "capability-map.xlsx")
    wb.save(out)
    print(f"  wrote {out}")


# ---------- 5) vendor-rfp-template.docx ----------

def make_vendor_rfp_template():
    doc = Document()
    _docx_set_margins(doc)
    _docx_title(
        doc,
        "AI 벤더 평가 RFP 템플릿",
        "문서번호: [RFP-AI-001]    ·    작성일: [YYYY-MM-DD]    ·    소관부서: [정보보호팀 / 구매팀 / AX TF]",
    )

    # 1. 목적과 범위
    doc.add_heading("1. 목적과 범위", level=1)
    meta_tbl = doc.add_table(rows=3, cols=2)
    meta_tbl.style = "Light Grid Accent 1"
    meta_rows = [
        ("회사명", "[ 회사명 ]"),
        ("작성일", "[ YYYY-MM-DD ]"),
        ("평가 대상 벤더 (3-5개)", "[ 예: Anthropic, OpenAI, Google, Microsoft, 사내 LLM ]"),
    ]
    for i, (k, v) in enumerate(meta_rows):
        meta_tbl.rows[i].cells[0].text = k
        meta_tbl.rows[i].cells[1].text = v
        for p in meta_tbl.rows[i].cells[0].paragraphs:
            for r in p.runs:
                r.bold = True
    doc.add_paragraph()
    doc.add_paragraph(
        "본 RFP는 사내 AI 도구 도입을 위한 표준 평가 절차입니다. "
        "최소 3개 이상의 벤더를 동일 기준으로 비교하며, 결정 매트릭스에 따라 메인·보조 벤더를 선정합니다."
    )

    # 2. 평가 기준
    doc.add_heading("2. 평가 기준", level=1)
    criteria = [
        ("기능", [
            "맥락 길이 (max context window)",
            "도구 사용 (function calling / MCP)",
            "이미지 입력·생성 지원",
            "코드 생성·실행 지원",
            "한국어 품질 (벤치마크 또는 사내 테스트)",
        ]),
        ("보안", [
            "데이터 학습·보존 정책 (zero-retention 옵션)",
            "한국 데이터센터 또는 리전 옵션",
            "SSO / SAML / SCIM 지원",
            "감사 로그 (audit log) 제공",
        ]),
        ("통합", [
            "REST API 안정성",
            "MCP 또는 동등 프로토콜",
            "Office / Workspace 통합",
        ]),
        ("가격", [
            "라이선스 모델 (시트 / 토큰 / 하이브리드)",
            "TCO (12개월 총비용)",
            "볼륨 할인 / 멀티이어 약정",
        ]),
        ("지원", [
            "한국어 기술 지원",
            "SLA (uptime / 응답 시간)",
            "온보딩·교육 지원",
        ]),
    ]
    for cat, items in criteria:
        doc.add_paragraph(cat, style="Intense Quote")
        for it in items:
            p = doc.add_paragraph(it, style="List Bullet")

    # 3. 벤더별 점수표
    doc.add_heading("3. 벤더별 점수표 (1-5점)", level=1)
    doc.add_paragraph(
        "각 항목 1-5점으로 평가하고 가중치를 곱해 합계를 계산합니다. "
        "가중치 합은 100점이며, 회사 우선순위에 따라 조정합니다."
    )

    score_tbl = doc.add_table(rows=7, cols=6)
    score_tbl.style = "Light Grid Accent 1"
    hdr = score_tbl.rows[0].cells
    for i, h in enumerate(["카테고리", "가중치", "벤더 A", "벤더 B", "벤더 C", "비고"]):
        hdr[i].text = h
    _bold_header_row(hdr)
    score_rows = [
        ("기능", "25", "", "", "", ""),
        ("보안", "30", "", "", "", ""),
        ("통합", "15", "", "", "", ""),
        ("가격", "20", "", "", "", ""),
        ("지원", "10", "", "", "", ""),
        ("가중 합계 (100점 환산)", "100", "", "", "", ""),
    ]
    for i, row in enumerate(score_rows, start=1):
        for j, v in enumerate(row):
            score_tbl.rows[i].cells[j].text = v

    doc.add_paragraph()

    # 4. 참조 고객 인터뷰 질문
    doc.add_heading("4. 참조 고객 인터뷰 질문", level=1)
    questions = [
        "도입 후 6개월 시점에 가장 만족스러운 점과 가장 아쉬운 점은 무엇입니까?",
        "사고(데이터 노출·잘못된 출력)가 발생한 적이 있다면 벤더의 대응은 어땠습니까?",
        "초기 약속한 SLA·기능 중 실제로 차이가 컸던 부분이 있습니까?",
        "라이선스·가격 협상에서 협상이 가능했던 항목과 불가능했던 항목은 무엇입니까?",
        "다시 도입을 결정한다면 같은 벤더를 선택하시겠습니까? 그 이유는 무엇입니까?",
    ]
    for q in questions:
        doc.add_paragraph(q, style="List Number")

    # 5. 결정 매트릭스
    doc.add_heading("5. 결정 매트릭스", level=1)
    doc.add_paragraph(
        "단일 벤더에 모든 것을 거는 것은 위험합니다. 메인 벤더 1개 + 보조 벤더 1개 조합을 권장합니다."
    )

    dec_tbl = doc.add_table(rows=4, cols=3)
    dec_tbl.style = "Light Grid Accent 1"
    hdr = dec_tbl.rows[0].cells
    for i, h in enumerate(["역할", "추천 벤더", "선정 사유"]):
        hdr[i].text = h
    _bold_header_row(hdr)
    dec_rows = [
        ("메인 벤더 (전사 표준)", "[ 가중 합계 1위 ]", "[ 기능·보안·통합 종합 점수 가장 높음 ]"),
        ("보조 벤더 (특수 케이스)", "[ 특정 카테고리 1위 ]", "[ 예: 한국어 1위 또는 가격 1위 ]"),
        ("재평가 시점", "[ 6개월 후 또는 사고 발생 시 ]", "[ SLA 위반·기능 격차 발생 시 즉시 ]"),
    ]
    for i, row in enumerate(dec_rows, start=1):
        for j, v in enumerate(row):
            dec_tbl.rows[i].cells[j].text = v

    _add_page_number_footer(doc, FOOTER)

    out = os.path.join(OUTDIR, "vendor-rfp-template.docx")
    doc.save(out)
    print(f"  wrote {out}")


# ---------- 6) architecture-patterns.pptx ----------

def make_architecture_patterns():
    prs = Presentation()
    prs.slide_width = Pin(13.333)
    prs.slide_height = Pin(7.5)

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Title
    title_box = slide.shapes.add_textbox(Pin(0.5), Pin(0.3), Pin(12.3), Pin(0.7))
    p = title_box.text_frame.paragraphs[0]
    p.text = "4가지 AI 아키텍처 패턴"
    p.font.size = PPt(26)
    p.font.bold = True
    p.font.color.rgb = PColor(0x11, 0x18, 0x27)
    p.alignment = PP_ALIGN.LEFT

    sub_box = slide.shapes.add_textbox(Pin(0.5), Pin(1.0), Pin(12.3), Pin(0.4))
    sp = sub_box.text_frame.paragraphs[0]
    sp.text = "통제·자동화 수준에 따라 패턴을 조합합니다. 단일 패턴만 쓰는 회사는 거의 없습니다."
    sp.font.size = PPt(12)
    sp.font.color.rgb = PColor(0x6B, 0x72, 0x80)

    # 2x2 grid
    ax, ay = 0.7, 1.6
    aw, ah = 11.93, 5.0
    cw, ch = aw / 2, ah / 2

    # (col, row, fill, title, desc, fit, risk)
    patterns = [
        (0, 0, (0xFE, 0xF3, 0xC7),
         "API only",
         "개별 사용자가 도구를 직접 사용 (Claude, ChatGPT 웹·앱)",
         "적합: 소규모 파일럿, 개인 생산성 도입 초기",
         "위험: 통제 약함 — 데이터 노출·도구 난립 가능"),
        (1, 0, (0xD1, 0xFA, 0xE5),
         "플러그인",
         "Office·Workspace에 임베디드 (M365 Copilot, Workspace AI)",
         "적합: 문서·이메일 중심 업무 50%+",
         "위험: 벤더 lock-in, 라이선스 비용 증가"),
        (0, 1, (0xDB, 0xEA, 0xFE),
         "코파일럿",
         "IDE·CRM·툴 내부에 직접 통합 (Cursor, Salesforce Einstein)",
         "적합: 특정 워크플로우 강화 (개발·영업)",
         "위험: 비용 상승, 도구별 라이선스 분산"),
        (1, 1, (0xFE, 0xE2, 0xE2),
         "에이전트",
         "자동 실행 워크플로우 (Claude Agent, AutoGen 등)",
         "적합: 반복적·규칙 명확한 백오피스 자동화",
         "위험: 검증 어려움, HITL 부재 시 사고 확산"),
    ]

    for col, row, rgb, title, desc, fit, risk in patterns:
        x = ax + col * cw
        y = ay + row * ch
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pin(x + 0.1), Pin(y + 0.1),
                                     Pin(cw - 0.2), Pin(ch - 0.2))
        shp.fill.solid()
        shp.fill.fore_color.rgb = PColor(*rgb)
        shp.line.color.rgb = PColor(0xFF, 0xFF, 0xFF)
        shp.line.width = PPt(2)
        tf = shp.text_frame
        tf.margin_left = Pin(0.25)
        tf.margin_top = Pin(0.2)
        tf.margin_right = Pin(0.25)
        tf.margin_bottom = Pin(0.2)
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = PPt(22)
        p1.font.bold = True
        p1.font.color.rgb = PColor(0x11, 0x18, 0x27)
        p1.alignment = PP_ALIGN.LEFT

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = PPt(12)
        p2.font.color.rgb = PColor(0x37, 0x41, 0x51)
        p2.space_before = PPt(6)

        p3 = tf.add_paragraph()
        p3.text = fit
        p3.font.size = PPt(11)
        p3.font.color.rgb = PColor(0x05, 0x73, 0x4A)
        p3.font.bold = True
        p3.space_before = PPt(8)

        p4 = tf.add_paragraph()
        p4.text = risk
        p4.font.size = PPt(11)
        p4.font.color.rgb = PColor(0x9B, 0x1C, 0x1C)
        p4.font.bold = True
        p4.space_before = PPt(3)

    # Note
    note_box = slide.shapes.add_textbox(Pin(0.5), Pin(6.65), Pin(12.3), Pin(0.4))
    np_ = note_box.text_frame.paragraphs[0]
    np_.text = "실제 회사는 보통 2-3개 패턴을 동시에 운영합니다 (예: 플러그인 + 코파일럿 + 일부 에이전트)."
    np_.font.size = PPt(10)
    np_.font.italic = True
    np_.font.color.rgb = PColor(0x6B, 0x72, 0x80)

    # Footer
    footer_box = slide.shapes.add_textbox(Pin(0.5), Pin(7.1), Pin(12.3), Pin(0.3))
    fp = footer_box.text_frame.paragraphs[0]
    fp.text = FOOTER
    fp.font.size = PPt(9)
    fp.font.italic = True
    fp.font.color.rgb = PColor(0x6B, 0x72, 0x80)
    fp.alignment = PP_ALIGN.RIGHT

    out = os.path.join(OUTDIR, "architecture-patterns.pptx")
    prs.save(out)
    print(f"  wrote {out}")


# ---------- 7) tco-comparison.xlsx ----------

def make_tco_comparison():
    wb = Workbook()
    ws = wb.active
    ws.title = "TCO 12개월 비교"

    set_col_widths(ws, [30, 16, 16, 16, 16, 16])

    put(ws, "A1", "TCO 12개월 비교 (벤더별)",
        Font(name="Inter", size=16, bold=True, color="111827"))
    ws.merge_cells("A1:F1")
    ws.row_dimensions[1].height = 26

    # Meta
    put(ws, "A3", "회사명", SUBHEADER_FONT, SUBHEADER_FILL, BORDER)
    put(ws, "B3", "", BODY_FONT, None, BORDER)
    ws.merge_cells("B3:C3")
    put(ws, "D3", "평가 기간", SUBHEADER_FONT, SUBHEADER_FILL, BORDER)
    put(ws, "E3", "", BODY_FONT, None, BORDER)
    ws.merge_cells("E3:F3")

    # Header row
    center = Alignment(horizontal="center", vertical="center", wrap_text=True)
    left = Alignment(horizontal="left", vertical="center", wrap_text=True)

    headers = ["항목 (만원)", "Claude", "GPT", "Gemini", "Copilot", "사내 LLM"]
    for i, h in enumerate(headers):
        put(ws, f"{get_column_letter(i+1)}5", h, HEADER_FONT, HEADER_FILL, BORDER, center)
    ws.row_dimensions[5].height = 24

    rows = [
        "라이선스 (연)",
        "통합 개발 (1회)",
        "운영 (연)",
        "교육 (1회)",
        "보안·감사 (연)",
        "사고 대응 적립금 (연)",
    ]

    start = 6
    for i, label in enumerate(rows):
        r = start + i
        put(ws, f"A{r}", label, BODY_FONT, None, BORDER, left)
        for c in range(2, 7):
            col = get_column_letter(c)
            put(ws, f"{col}{r}", "", BODY_FONT, None, BORDER, center)

    # 합계 행 — 자동 합산
    total_row = start + len(rows)
    put(ws, f"A{total_row}", "합계", SUBHEADER_FONT, SUBHEADER_FILL, BORDER, left)
    for c in range(2, 7):
        col = get_column_letter(c)
        formula = f"=SUM({col}{start}:{col}{total_row-1})"
        put(ws, f"{col}{total_row}", formula, BODY_FONT, ACCENT_FILL, BORDER, center)

    # 가이드 안내
    note_row = total_row + 2
    put(ws, f"A{note_row}",
        "노랑 셀(합계)은 자동 계산됩니다. 각 비용은 만원 단위로 입력하세요.",
        HINT_FONT)
    ws.merge_cells(f"A{note_row}:F{note_row}")

    note_row2 = note_row + 1
    put(ws, f"A{note_row2}",
        "참고: 사내 LLM은 운영 비용이 크지만 데이터 등급 Lv3-Lv4 처리가 가능합니다. "
        "외부 SaaS는 데이터 등급에 따라 사용 제약이 있으므로 단순 가격 비교만으로 결정하지 마세요.",
        HINT_FONT, align=Alignment(wrap_text=True, vertical="top"))
    ws.merge_cells(f"A{note_row2}:F{note_row2}")
    ws.row_dimensions[note_row2].height = 36

    footer_row = note_row2 + 2
    put(ws, f"A{footer_row}", FOOTER, FOOTER_FONT)
    ws.merge_cells(f"A{footer_row}:F{footer_row}")

    out = os.path.join(OUTDIR, "tco-comparison.xlsx")
    wb.save(out)
    print(f"  wrote {out}")


# ---------- main ----------

if __name__ == "__main__":
    os.makedirs(OUTDIR, exist_ok=True)
    make_incident_scenarios()
    make_job_redesign_workbook()
    make_operating_models()
    make_capability_map()
    make_vendor_rfp_template()
    make_architecture_patterns()
    make_tco_comparison()
    print("Generated 7 P2-B assets")
