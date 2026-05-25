"""
Build P3 (Execution track) download assets for AX Guide for Enterprise.

Outputs (in same directory as this script):
- prompt-pack.md
- claude-md-template.md
- internal-glossary-template.md
- internal-guide-template.md
- gemini-md-template.md
- gpt-project-instructions-template.md
- tool-pricing-comparison.xlsx
- cases-30-list.xlsx
- case-collection-form.docx
- internal-guide-template.docx
- onboarding-1pagers.docx
- monthly-retrospective.docx
- case-presentation.pptx
- tool-decision-tree.docx        (PDF substitute since reportlab not installed)

Run from repo root:
    python3 assets/downloads/_build_assets_p3.py
"""

import os

from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from pptx import Presentation
from pptx.util import Inches as Pin, Pt as PPt
from pptx.dml.color import RGBColor as PColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTDIR = SCRIPT_DIR
FOOTER = "AX Guide for Enterprise · v0.2 / jade@linercorp.com"


# ---------- common style helpers ----------

THIN = Side(border_style="thin", color="CCCCCC")
BORDER = Border(left=THIN, right=THIN, top=THIN, bottom=THIN)

HEADER_FILL = PatternFill("solid", fgColor="1F2937")
SUBHEADER_FILL = PatternFill("solid", fgColor="E5E7EB")
ACCENT_FILL = PatternFill("solid", fgColor="FEF3C7")
GOOD_FILL = PatternFill("solid", fgColor="D1FAE5")
WARN_FILL = PatternFill("solid", fgColor="FEE2E2")

HEADER_FONT = Font(name="Inter", size=12, bold=True, color="FFFFFF")
SUBHEADER_FONT = Font(name="Inter", size=11, bold=True, color="111827")
BODY_FONT = Font(name="Inter", size=11, color="111827")
HINT_FONT = Font(name="Inter", size=10, italic=True, color="6B7280")
FOOTER_FONT = Font(name="Inter", size=9, italic=True, color="6B7280")


def set_col_widths(ws, widths):
    for i, w in enumerate(widths, start=1):
        ws.column_dimensions[get_column_letter(i)].width = w


def put(ws, cell, value, font=None, fill=None, border=None, align=None):
    c = ws[cell]
    c.value = value
    if font:
        c.font = font
    if fill:
        c.fill = fill
    if border:
        c.border = border
    if align:
        c.alignment = align


def _add_page_number_footer(doc, footer_text):
    section = doc.sections[0]
    footer = section.footer
    p = footer.paragraphs[0]
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    run = p.add_run(f"{footer_text}    ·    ")
    run.font.size = Pt(9)
    run.font.italic = True
    run.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

    fld_run = p.add_run()
    fld_run.font.size = Pt(9)
    fld_run.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

    fldChar1 = OxmlElement("w:fldChar")
    fldChar1.set(qn("w:fldCharType"), "begin")
    instrText = OxmlElement("w:instrText")
    instrText.set(qn("xml:space"), "preserve")
    instrText.text = "PAGE"
    fldChar2 = OxmlElement("w:fldChar")
    fldChar2.set(qn("w:fldCharType"), "end")
    fld_run._r.append(fldChar1)
    fld_run._r.append(instrText)
    fld_run._r.append(fldChar2)

    fld_run.add_text(" / ")

    fldChar3 = OxmlElement("w:fldChar")
    fldChar3.set(qn("w:fldCharType"), "begin")
    instrText2 = OxmlElement("w:instrText")
    instrText2.set(qn("xml:space"), "preserve")
    instrText2.text = "NUMPAGES"
    fldChar4 = OxmlElement("w:fldChar")
    fldChar4.set(qn("w:fldCharType"), "end")
    fld_run._r.append(fldChar3)
    fld_run._r.append(instrText2)
    fld_run._r.append(fldChar4)


def _docx_set_margins(doc, inches=1.0):
    for section in doc.sections:
        section.top_margin = Inches(inches)
        section.bottom_margin = Inches(inches)
        section.left_margin = Inches(inches)
        section.right_margin = Inches(inches)


def _docx_title(doc, title_text, meta_text=None):
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run(title_text)
    run.bold = True
    run.font.size = Pt(22)
    run.font.color.rgb = RGBColor(0x11, 0x18, 0x27)

    if meta_text:
        meta = doc.add_paragraph()
        meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
        mr = meta.add_run(meta_text)
        mr.font.size = Pt(10)
        mr.font.italic = True
        mr.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
    doc.add_paragraph()


def _bold_header_row(row_cells):
    for cell in row_cells:
        for p in cell.paragraphs:
            for r in p.runs:
                r.bold = True


# ===================== MARKDOWN FILES =====================

def write_md(name, content):
    path = os.path.join(OUTDIR, name)
    with open(path, "w", encoding="utf-8") as f:
        f.write(content)
    print(f"  wrote {name} ({len(content)} bytes)")


PROMPT_PACK_MD = """# AX Guide · 프롬프트 · 하네스 패턴팩 v0.2

> 15개 패턴 (기본 5 + 하네스 5 + 구조 5). 각 패턴마다 한 줄 정의 + 예시 1개.
>
> 출처: AX Guide for Enterprise — https://github.com/[org]/ax-guide

## A. 기본 패턴 5종

### A1. 역할 부여
**정의:** 모델에 명시적 역할을 주면 출력 톤·관점이 달라진다.
**예시:**
```
당신은 시니어 카피라이터입니다. 다음 카피를 B2B SaaS 마케팅 톤으로 다시 써주세요:
[원본 카피]
```

### A2. 단계 분리
**정의:** 한 프롬프트에 여러 단계를 묶지 말고, 명시적으로 단계화한다.
**예시:**
```
다음 작업을 두 단계로 수행해 주세요:
1단계: 회의록을 한국어 5문장으로 요약
2단계: 그 요약을 기반으로 임원 보고용 3줄 메모 작성
[회의록]
```

### A3. 예시 제공 (Few-shot)
**정의:** 원하는 출력 형식의 예시 2-3개를 먼저 주면 정확도가 크게 올라간다.
**예시:**
```
다음 형식으로 제품 설명을 작성해 주세요.

예시 1:
- 한 줄 가치: 회의록 자동 요약
- 핵심 기능: 5분 회의 → 30초 요약
- 차별점: 한국어 도메인 특화

예시 2:
- 한 줄 가치: 견적서 자동 작성
- 핵심 기능: 항목 입력 → 즉시 PDF
- 차별점: 사내 ERP 연동

이제 다음 제품을 같은 형식으로:
[제품 정보]
```

### A4. 검증 요청
**정의:** 출력 직후 모델 본인이 자기 출력을 검증하게 하면 사실 오류·논리 모순을 잡아낸다.
**예시:**
```
다음 보고서를 작성하고, 작성 직후 다음 3가지를 자체 검증해 주세요:
(1) 사실 오류 가능성이 있는 문장 3개
(2) 논리적 비약이 있는 단락
(3) 출처가 명시되지 않은 주장

[과제]
```

### A5. 출력 형식 지정
**정의:** Markdown 표·JSON·체크리스트 등 출력 형식을 미리 지정한다.
**예시:**
```
다음 데이터를 분석해서 아래 JSON 형식으로만 응답해 주세요. 자연어 설명은 제외.

{
  "headline": "...",
  "key_findings": ["...", "...", "..."],
  "risks": ["...", "..."],
  "next_actions": ["...", "...", "..."]
}

[데이터]
```

## B. 하네스 패턴 5종 (AI ROASTING 고유)

### B1. 5-Color Harness
**정의:** BLACK(수행자) / RED(이성 비평) / SILVER(분야 전문가) / BLUE(공감) / GOLD(독자) — 5 페르소나로 같은 산출물을 5번 평가.
**예시:**
```
다음 보고서 초안에 대해 5-Color Harness 평가를 해주세요.
- BLACK: 수행자 관점에서 빠진 단계
- RED: 이성 관점에서 논리·근거 약점
- SILVER: 도메인 전문가 관점에서 디테일 오류
- BLUE: 공감 관점에서 독자 감정 흐름
- GOLD: 최종 독자 관점에서 첫 반응 + 다음 액션

각 색마다 3가지씩 코멘트.

[보고서 초안]
```

### B2. 멀티 페르소나 토론
**정의:** 3-5명의 페르소나가 같은 안건을 각자 입장에서 발언. 합의·이견·잔여 리스크 정리.
**예시:**
```
다음 의사결정 안건에 대해 4인 페르소나 토론을 시뮬레이션:
- CFO: 비용 관점
- CTO: 기술 가능성
- CCO: 고객 영향
- CHRO: 인재 영향

각자 3분 발언 → 30초 합의 시도 → 잔여 이견 + 결정 제안.

[안건]
```

### B3. 7·9.5·10점 등급 기준
**정의:** 산출물마다 목표 등급을 명시. 7점=빠른 실험, 9.5점=실사용 제품, 10점=대표작/외부 발표.
**예시:**
```
다음 산출물을 9.5점 기준으로 작성해 주세요. 9.5점 기준:
- 외부에서 그대로 인용 가능
- 사람이 윤문하면 10점 도달
- 시간 절감보다 품질이 우선

[과제]
```

### B4. 자기검증 루프
**정의:** 출력 → 자기 평가 → 약점 도출 → 재작성. 3-pass 권장.
**예시:**
```
다음 작업을 3-pass로 수행:
Pass 1: 초안 작성
Pass 2: 자기 평가 (강점·약점 3가지씩)
Pass 3: 약점 보강한 재작성

각 pass 결과를 모두 보여주세요.

[과제]
```

### B5. 외부 송출 게이트
**정의:** 외부 송출 자산은 RGSB(Reason/Domain/Empathy/Reader) 4관점 통과 필수.
**예시:**
```
다음 자산을 외부 송출 전 RGSB 게이트로 점검:
- Reason: 논리·근거 약점 3가지
- Domain: 도메인 디테일 오류
- Empathy: 독자 감정 흐름
- Reader: 최종 독자 1차 반응 + 다음 액션

각 관점에서 PASS/FAIL + 수정 권고.

[자산]
```

## C. 반복 가능한 구조 패턴 5종

### C1. 시스템 프롬프트
**정의:** 매번 같은 컨텍스트를 반복하지 말고 시스템 프롬프트(또는 프로젝트 지침)에 한 번 박는다.
**예시:**
```
[시스템 프롬프트로 한 번만 설정]

당신은 [회사명]의 마케팅팀 AI 비서다.
- 우리 고객은 B2B SaaS의 CFO/COO다.
- 톤: 명료하고 단정한 한국어, 과장어 금지.
- 출력 순서: 핵심 판단 → 이유 → 실행안 → 다음 액션.
- 모르는 사실은 추측 대신 명시적으로 "확인 필요"라고 말한다.
```

### C2. 프로젝트 컨텍스트
**정의:** Claude Project / ChatGPT Project 단위로 컨텍스트 파일을 묶어 항상 함께 로드되게 만든다.
**예시:**
```
[프로젝트에 첨부할 파일들]
- glossary.md (사내 어휘집)
- brand-voice.md (브랜드 톤)
- past-cases.md (과거 케이스 5건)

[프로젝트 지침]
위 3개 파일을 항상 참조해서 답해 주세요.
새 산출물은 brand-voice.md 톤에 맞춰 작성.
```

### C3. 스킬 · MCP
**정의:** 자주 쓰는 워크플로우는 스킬·MCP 서버로 자동화. 단발 프롬프트 → 재사용 가능한 도구로 승격.
**예시:**
```
[스킬 정의 예시 - SKILL.md]

name: weekly-status
description: 주간 업무 보고서 자동 생성

trigger: 사용자가 "주간 보고", "이번 주 정리", "weekly status"를 언급할 때

input:
- 이번 주 슬랙 메시지 (선택)
- 이번 주 캘린더 이벤트 (선택)

output: 5-section 보고서 (성과·진행중·블록·다음주·리스크)
```

### C4. CLAUDE.md / GEMINI.md / GPT 프로젝트 지침
**정의:** Claude Code·Gemini·Custom GPT 각각의 컨텍스트 파일에 회사 컨벤션을 박는다.
**예시:**
```
[CLAUDE.md 예시 - 코딩 프로젝트 루트에 둠]

# CLAUDE.md

## 컨벤션
- 언어: TypeScript + Next.js
- 들여쓰기: 2 spaces
- 커밋: Conventional Commits

## 절대 금지
- 프로덕션 DB 직접 쿼리
- 외부 LLM에 고객 데이터 입력

## 기본 출력
1. 핵심 판단 2. 이유 3. 실행안 4. 다음 액션
```

### C5. 사내 어휘집
**정의:** 회사 내부 어휘 사전. LLM이 회사 어휘로 답하게 만든다. RAG에 가장 먼저 들어갈 자료.
**예시:**
```
다음 사내 어휘집을 컨텍스트로 사용:
- 고객 = B2B SaaS 회사 의사결정자 (보통 CFO 또는 COO)
- 성과 = 월간 활성 사용자(MAU) + 매출 둘 다
- 핵심 가치 = 신뢰·속도·투명성

위 어휘를 일관되게 사용해서 다음 콘텐츠 작성:
[과제]
```

---

> 이 패턴팩은 분기 1회 업데이트됩니다.
> AX Guide for Enterprise · v0.2 / jade@linercorp.com
"""


CLAUDE_MD_TEMPLATE = """# CLAUDE.md

이 파일은 Claude Code가 이 프로젝트에서 작업할 때 참고하는 지침이다.

## 프로젝트 컨텍스트

- **회사/팀:** [회사명]
- **프로젝트 목적:** [한 줄 정의]
- **주요 사용자:** [누구]
- **현재 단계:** [기획 / 개발 / 운영]

## 코딩 컨벤션

- 언어/프레임워크: [예: Python 3.11 + FastAPI, TypeScript + Next.js]
- 들여쓰기: [2 spaces / 4 spaces / tab]
- 네이밍: [camelCase / snake_case]
- 주석: [필수 항목 / 영문 vs 한글]

## 작업 원칙

- TDD 적용 여부: [예/아니오, 적용 범위]
- 커밋 메시지 컨벤션: [Conventional Commits 사용 여부]
- PR 단위: [1 feature = 1 PR / multi-task PR 허용]

## 핵심 자산

- `docs/architecture.md` — 시스템 아키텍처 (참조 필수)
- `docs/glossary.md` — 사내 어휘집
- `scripts/setup.sh` — 로컬 환경 부트스트랩

## 절대 하지 말 것

- 프로덕션 DB에 직접 쿼리
- 비밀 키를 코드에 인라인 (사용 시 환경 변수)
- 외부 LLM에 고객 데이터 입력 (사내 LLM만 허용)
- 검증 없이 자동화 워크플로우 배포

## 기본 출력 패턴

1. 핵심 판단
2. 이유
3. 실행안
4. 남길 자산
5. 다음 액션

---

> AX Guide for Enterprise · v0.2 / 사내 사용자가 회사에 맞춰 채워 사용.
> 원본 패턴은 https://github.com/[org]/ax-guide
"""


INTERNAL_GLOSSARY_MD = """# 사내 어휘집 (Internal Glossary)

> AI 도구가 우리 회사 어휘로 답하게 만드는 RAG 첫 번째 자료. 시스템 프롬프트에 통째로 주입하거나 RAG 컨텍스트로 활용.

## 사용 방법

1. 빈 칸을 회사 사정에 맞게 채운다
2. CLAUDE.md / GPT 프로젝트 지침에 첨부하거나 사내 LLM RAG에 업로드
3. 분기 1회 갱신

## 어휘 정의

### 비즈니스 핵심

- **고객 (Customer)** — [구체적으로 누구. 예: "B2B SaaS 회사의 CFO"]
- **유저 (User)** — [고객 회사 내 실제 도구 사용자]
- **성과 (Outcome)** — [무엇을 성과로 보는가. 예: "MAU + 매출 둘 다"]
- **핵심 가치 (Core Value)** — [회사의 3가지 가치]

### 제품·서비스

- **메인 제품** — [정의 + 차별점]
- **부가 서비스** — [정의]
- **로드맵** — [현재 6개월 계획]

### 조직·역할

- **CEO/COO/CFO/CTO** — [현재 누구]
- **각 부서 리더** — [영업·마케팅·HR·재무·개발·디자인·고객지원 리더]
- **AX 챔피언** — [각 부서 AX 챔피언]

### 의사결정 용어

- **결정 등급** — [예: "임원 / 부서장 / 팀 / 개인"]
- **승인 프로세스** — [데이터 등급별 승인자]

### 도구 환경

- **사용 가능 AI 도구** — [Claude / GPT / Gemini / Copilot / 사내 LLM 중 어떤 것]
- **데이터 등급 분류** — [4등급 정의]

---

> AX Guide for Enterprise · v0.2 / 사내 어휘집 분기 1회 갱신.
"""


INTERNAL_GUIDE_MD = """# [회사명] AI 사용 가이드 v1.0

> 작성: [날짜] / 작성자: [AX 챔피언 이름] / 차기 갱신: [3개월 뒤]

## 1. 우리 회사의 AX 비전

[한 문장. 예: "고객의 신뢰를 지키면서 제품 개발·고객 지원·콘텐츠 제작을 AI로 재설계해서, 3년 뒤 1인당 매출 3배의 조직이 된다."]

## 2. 사용 가능한 AI 도구

| 도구 | 등급 | 사용 가능 데이터 | 비용 | 라이선스 |
|---|---|---|---|---|
| Claude (외부) | 메인 | 공개·내부 | 월 $20 | 개인 결제 |
| ChatGPT (외부) | 보조 | 공개·내부 | 월 $20 | 개인 결제 |
| Gemini Workspace | 메인 (Google 환경) | 공개·내부·기밀 | 회사 라이선스 | 회사 발급 |
| 사내 LLM | 기밀 데이터용 | 모든 등급 | 인프라 운영비 | 사내 운영 |

## 3. 데이터 등급별 사용 규칙

- **Lv1 공개:** 모든 도구 OK
- **Lv2 내부:** 외부 LLM은 가공 후 OK
- **Lv3 기밀:** 사내 LLM만 OK
- **Lv4 극비:** 사내 LLM + 임원 승인

## 4. 금지 행위

1. 극비 데이터(고객 개인정보·계약·재무 디테일)를 외부 LLM에 입력
2. AI 출력물을 검증 없이 외부 송출 (외부 송출 게이트 통과 필수)
3. 정책 미승인 무료 도구 사용

## 5. 외부 송출 게이트

외부에 나가는 모든 AI 협업 산출물은 다음 RGSB 4관점 통과 후 송출:
- **R**eason: 논리·근거 OK?
- **G** (Domain): 도메인 디테일 정확?
- **S** (Empathy): 독자 감정 흐름 OK?
- **B** (Reader): 최종 독자 첫 반응 OK?

## 6. 사고 발생 시

(1) 즉시 보고 → (2) 영향 평가 → (3) 외부 통보 (필요 시) → (4) 사후 분석 → (5) 재발 방지 룰 추가

## 7. 챔피언 명단

- [부서] [이름] — 슬랙 @[handle]

## 8. 분기 회고

- 분기말 1시간 워크숍
- 새 케이스 5개 공유
- 가이드 v1.1 / v1.2 ... 업데이트

---

> 본 가이드는 AX Guide for Enterprise (v0.2) 템플릿을 기반으로 작성되었습니다.
"""


GEMINI_MD_TEMPLATE = """# Gemini Gem · [프로젝트명] 지침

이 Gem은 Google Workspace 환경에서 [팀명] 업무를 보조한다.

## 컨텍스트

- **회사/팀:** [회사명]
- **프로젝트 목적:** [한 줄 정의]
- **주요 사용자:** [누구]
- **현재 단계:** [기획 / 운영 / 확산]

## Workspace 특화 패턴

- **Docs/Sheets/Slides 통합 활용** — Drive 내 자료를 직접 참조해 답한다. 새 산출물도 동일 포맷으로.
- **Drive 자료 직접 인용** — 답변 시 출처 Doc/Sheet 링크를 명시한다.
- **Calendar/Gmail 연동 워크플로우** — 일정·메일 컨텍스트가 필요한 작업은 해당 도구의 컨텍스트를 활용한다.

## 코딩/문서 컨벤션

- 작성 언어: 한국어 기본, 영문 인용 시 원문 병기
- 톤: 명료하고 단정한 한국어, 과장어 금지
- 출력 순서: 핵심 판단 → 이유 → 실행안 → 다음 액션

## 절대 하지 말 것

- 극비(Lv4) 데이터를 외부 Gem에 입력
- 검증 없이 외부 메일/문서 자동 송출
- Drive 자료의 출처 누락

## 기본 출력 패턴

1. 핵심 판단
2. 이유
3. 실행안
4. 남길 자산
5. 다음 액션

---

> AX Guide for Enterprise · v0.2 / Gemini Gem 또는 Workspace 프로젝트용 템플릿.
"""


GPT_PROJECT_INSTRUCTIONS_MD = """# Custom GPT · [프로젝트명]

당신은 [팀명]의 [역할] AI 비서입니다.

## 컨텍스트 요약

- **회사/팀:** [회사명]
- **프로젝트 목적:** [한 줄 정의]
- **주요 사용자:** [누구]
- **현재 단계:** [기획 / 운영 / 확산]

## 응답 톤

- 명확하고 단정한 한국어
- 결정 → 이유 → 실행안 → 다음 액션 순서
- 추측보다 명시적 질문
- 과장어("혁신적", "획기적") 금지

## 도구 사용

- **DALL-E:** 콘텐츠 시각화 요청 시
- **Code Interpreter:** 데이터 분석·계산
- **Web Browsing:** 최신 정보 필요 시 (단 출처 명시)
- **첨부 파일(Knowledge):** 사내 어휘집·과거 케이스·브랜드 가이드 첨부

## 절대 하지 말 것

- 극비(Lv4) 데이터 입력 금지 (외부 GPT 사용 시)
- 검증 없이 외부 송출용 답변 작성
- 출처 없는 단정적 사실 진술

## 기본 출력 패턴

1. 핵심 판단
2. 이유
3. 실행안
4. 남길 자산
5. 다음 액션

---

> AX Guide for Enterprise · v0.2 / Custom GPT 프로젝트 지침 템플릿.
"""


def make_markdown_files():
    write_md("prompt-pack.md", PROMPT_PACK_MD)
    write_md("claude-md-template.md", CLAUDE_MD_TEMPLATE)
    write_md("internal-glossary-template.md", INTERNAL_GLOSSARY_MD)
    write_md("internal-guide-template.md", INTERNAL_GUIDE_MD)
    write_md("gemini-md-template.md", GEMINI_MD_TEMPLATE)
    write_md("gpt-project-instructions-template.md", GPT_PROJECT_INSTRUCTIONS_MD)


# ===================== EXCEL FILES =====================

def make_tool_pricing_comparison():
    wb = Workbook()
    ws = wb.active
    ws.title = "도구 가격 비교"

    set_col_widths(ws, [22, 18, 22, 22, 30, 22])

    # Title
    ws.merge_cells("A1:F1")
    put(ws, "A1", "AI 도구 가격·정책 비교 (2026년 5월 기준)",
        font=Font(name="Inter", size=16, bold=True, color="111827"),
        align=Alignment(horizontal="center", vertical="center"))
    ws.row_dimensions[1].height = 30

    ws.merge_cells("A2:F2")
    put(ws, "A2", "단가는 USD 월정액 / 변동 가능 — 분기 1회 갱신",
        font=HINT_FONT,
        align=Alignment(horizontal="center"))

    # Header row
    headers = ["도구", "개인 (월, $)", "팀 (사용자당, $/월)", "엔터프라이즈", "데이터 학습 정책", "한국 데이터센터"]
    for i, h in enumerate(headers, start=1):
        put(ws, f"{get_column_letter(i)}4", h,
            font=HEADER_FONT, fill=HEADER_FILL, border=BORDER,
            align=Alignment(horizontal="center", vertical="center"))
    ws.row_dimensions[4].height = 24

    data = [
        ("Claude", "20", "25-30", "별도 협의", "학습 안 함 (기본)", "X"),
        ("ChatGPT", "20", "25", "60+", "학습 안 함 (Team+ 옵션)", "X"),
        ("Gemini", "20", "22 (Workspace 통합)", "별도", "Workspace 라이선스 보호", "O (asia-northeast3)"),
        ("Microsoft Copilot", "-", "30", "30/사용자", "M365 라이선스 보호", "O (Korea Central)"),
        ("Perplexity", "20", "40", "별도", "학습 안 함 (Pro+)", "X"),
        ("NotebookLM", "무료", "무료", "무료 (베타)", "학습 안 함", "X"),
        ("Cursor", "20", "20", "40", "학습 안 함 (Privacy Mode)", "X"),
        ("Claude Code", "API 종량", "API 종량", "API 종량", "학습 안 함", "X"),
        ("사내 LLM", "-", "-", "인프라 운영비", "100% 내부", "O (사내 서버)"),
        ("API 직접", "-", "-", "토큰당", "학습 안 함 (API)", "벤더별"),
    ]

    for r, row in enumerate(data, start=5):
        for c, val in enumerate(row, start=1):
            font = SUBHEADER_FONT if c == 1 else BODY_FONT
            fill = SUBHEADER_FILL if c == 1 else None
            align = Alignment(horizontal="left", vertical="center", wrap_text=True) if c in (1, 5, 6) \
                else Alignment(horizontal="center", vertical="center", wrap_text=True)
            put(ws, f"{get_column_letter(c)}{r}", val,
                font=font, fill=fill, border=BORDER, align=align)
        ws.row_dimensions[r].height = 26

    # Footer
    footer_row = 5 + len(data) + 2
    ws.merge_cells(f"A{footer_row}:F{footer_row}")
    put(ws, f"A{footer_row}", FOOTER,
        font=FOOTER_FONT, align=Alignment(horizontal="center"))

    out = os.path.join(OUTDIR, "tool-pricing-comparison.xlsx")
    wb.save(out)
    print(f"  wrote tool-pricing-comparison.xlsx")


CASES_25 = [
    # (카테고리, 케이스명, 직군, 도구, Before, After, 시간절감)
    ("경영기획", "월간 임원 보고서 자동 초안", "경영기획", "Claude + 사내 어휘집",
     "데이터 수집·차트·문장 6시간", "데이터 입력 → 초안 30분 + 윤문 1시간", "60%"),
    ("경영기획", "전략 메모 다중 페르소나 검토", "경영기획", "Claude + 5-Color Harness",
     "임원 회람 후 1주일 피드백", "5인 페르소나 평가 1시간", "85%"),
    ("경영기획", "이사회 자료 사전 QA", "경영기획", "Claude (외부 송출 게이트)",
     "팀장이 4시간 검토", "RGSB 게이트 1시간 + 팀장 30분", "70%"),
    ("영업", "고객 메일 톤 맞춤 작성", "영업", "ChatGPT + 사내 어휘집",
     "1통당 20분", "1통당 5분 (윤문 포함)", "75%"),
    ("영업", "RFP 응답서 초안 자동화", "영업", "Claude Project + 과거 RFP",
     "RFP 1건당 2-3일", "초안 4시간 + 윤문 1일", "60%"),
    ("영업", "고객 미팅 사전 리서치", "영업", "Perplexity + ChatGPT",
     "미팅당 1.5시간", "미팅당 20분", "75%"),
    ("마케팅", "SNS 카피 5종 동시 생성", "마케팅", "ChatGPT + 브랜드 가이드",
     "카피라이터 반나절", "30분 생성 + 1시간 윤문", "70%"),
    ("마케팅", "이메일 캠페인 A/B 안 작성", "마케팅", "Claude + 과거 성과 데이터",
     "캠페인당 4시간", "A/B 6안 1시간 + 검증 30분", "65%"),
    ("마케팅", "블로그 SEO 초안 작성", "마케팅", "Claude + 키워드 리서치",
     "글 1편당 1일", "초안 2시간 + 윤문 2시간", "50%"),
    ("HR", "JD 채용 공고 자동 작성", "HR", "Claude + 회사 JD 템플릿",
     "JD 1건 3시간", "JD 1건 45분", "75%"),
    ("HR", "면접 질문 셋 자동 생성", "HR", "Claude + 직무 요건",
     "면접관별 1시간", "5세트 30분 + 면접관 리뷰 15분", "70%"),
    ("HR", "온보딩 안내문 개인 맞춤", "HR", "ChatGPT + 신입 프로필",
     "1인당 1시간", "1인당 10분", "80%"),
    ("재무", "월 결산 변동 분석 요약", "재무·회계", "Claude + 엑셀 데이터",
     "분석가 1일", "분석가 2시간 + 검증 1시간", "60%"),
    ("재무", "예산 차이 설명문 자동 작성", "재무·회계", "Claude + 사내 어휘집",
     "부서별 1시간", "전체 30분 자동 + 검증 30분", "75%"),
    ("재무", "투자 검토 메모 1차 초안", "재무·회계", "Claude + 과거 메모",
     "메모 1건 2일", "초안 3시간 + 검토 반나절", "65%"),
    ("개발", "PR 리뷰 자동 보조", "개발", "Claude Code + 사내 컨벤션",
     "리뷰어 1시간", "1차 리뷰 15분 + 사람 30분", "55%"),
    ("개발", "테스트 코드 자동 생성", "개발", "Cursor + 기존 코드",
     "함수당 30분", "함수당 5분", "85%"),
    ("개발", "버그 리포트 → 재현 코드", "개발", "Claude Code",
     "재현 1시간", "재현 15분", "75%"),
    ("디자인", "디자인 시안 텍스트 자동 채우기", "디자인", "ChatGPT + 브랜드 톤",
     "시안당 1시간", "시안당 10분", "85%"),
    ("디자인", "사용자 인터뷰 인사이트 추출", "디자인", "Claude + 녹취록",
     "1시간 인터뷰 분석 2시간", "30분 분석 + 검토 30분", "65%"),
    ("디자인", "디자인 리뷰 코멘트 정리", "디자인", "Claude",
     "리뷰 후 1시간 정리", "10분 자동 + 디자이너 10분", "80%"),
    ("고객지원", "FAQ 답변 1차 자동화", "고객지원", "사내 LLM + FAQ DB",
     "건당 10-15분", "건당 1분 (1차) + 검토 5분", "60%"),
    ("고객지원", "고객 불만 톤 분석·우선순위", "고객지원", "Claude + 티켓 데이터",
     "일간 정리 2시간", "일간 정리 30분", "75%"),
    ("고객지원", "복잡 케이스 대응 시나리오 생성", "고객지원", "Claude + 과거 케이스",
     "시니어가 30분", "5분 자동 + 시니어 10분 검토", "70%"),
    ("고객지원", "다국어 응대 초안", "고객지원", "Claude + 사내 어휘집",
     "번역사 의뢰 1일", "초안 5분 + 윤문 30분", "85%"),
]


def make_cases_30_list():
    wb = Workbook()
    ws = wb.active
    ws.title = "사내 케이스 일람"

    set_col_widths(ws, [14, 30, 14, 26, 32, 32, 12])

    ws.merge_cells("A1:G1")
    put(ws, "A1", "AI 협업 사내 케이스 일람 (25선 + 추가 슬롯)",
        font=Font(name="Inter", size=16, bold=True, color="111827"),
        align=Alignment(horizontal="center", vertical="center"))
    ws.row_dimensions[1].height = 30

    ws.merge_cells("A2:G2")
    put(ws, "A2", "분기 1회 사내 컨테스트로 신규 케이스를 모읍니다.",
        font=HINT_FONT, align=Alignment(horizontal="center"))

    headers = ["카테고리", "케이스명", "직군", "도구", "Before", "After", "시간 절감"]
    for i, h in enumerate(headers, start=1):
        put(ws, f"{get_column_letter(i)}4", h,
            font=HEADER_FONT, fill=HEADER_FILL, border=BORDER,
            align=Alignment(horizontal="center", vertical="center"))
    ws.row_dimensions[4].height = 24

    for r, row in enumerate(CASES_25, start=5):
        for c, val in enumerate(row, start=1):
            align = Alignment(horizontal="left", vertical="center", wrap_text=True) if c in (2, 4, 5, 6) \
                else Alignment(horizontal="center", vertical="center", wrap_text=True)
            put(ws, f"{get_column_letter(c)}{r}", val,
                font=BODY_FONT, border=BORDER, align=align)
        ws.row_dimensions[r].height = 38

    # Empty slots (5 rows)
    base = 5 + len(CASES_25)
    for i in range(5):
        r = base + i
        put(ws, f"A{r}", "[사내 케이스 추가]",
            font=HINT_FONT, border=BORDER, fill=ACCENT_FILL,
            align=Alignment(horizontal="center", vertical="center"))
        for c in range(2, 8):
            put(ws, f"{get_column_letter(c)}{r}", "",
                font=BODY_FONT, border=BORDER, fill=ACCENT_FILL)
        ws.row_dimensions[r].height = 34

    footer_row = base + 5 + 2
    ws.merge_cells(f"A{footer_row}:G{footer_row}")
    put(ws, f"A{footer_row}", FOOTER,
        font=FOOTER_FONT, align=Alignment(horizontal="center"))

    out = os.path.join(OUTDIR, "cases-30-list.xlsx")
    wb.save(out)
    print(f"  wrote cases-30-list.xlsx")


# ===================== WORD FILES =====================

def make_case_collection_form():
    doc = Document()
    _docx_set_margins(doc)
    _docx_title(
        doc,
        "AI 협업 사례 수집 양식",
        "AX Guide for Enterprise · 분기 1회 사내 컨테스트용",
    )

    doc.add_heading("1. 목적", level=1)
    doc.add_paragraph(
        "AI 협업 사례를 사내에 공유하여 학습·확산을 가속화합니다. "
        "각 케이스는 한 장에 담길 수 있도록 압축하고, 분기 컨테스트에 제출됩니다."
    )

    doc.add_heading("2. 수집 양식 (1 케이스당 1장)", level=1)

    fields = [
        ("제목 (한 줄)", "예: 월간 임원 보고서 자동 초안 (60% 시간 절감)"),
        ("작성자 · 부서 · 날짜", "예: 김OO / 경영기획팀 / 2026-05-24"),
        ("Before — 현재 일하는 방식", "도구·시간·품질 관점에서 구체적으로. 예: 데이터 수집·차트·문장 작성 합 6시간."),
        ("도구 · 프롬프트 · 검증 방법", "도구·프롬프트 구조·검증 단계를 구체적으로. 예: Claude Project + 사내 어휘집 + 5-Color Harness 검증."),
        ("After — 바뀐 일하는 방식", "구체적인 시간·품질 변화. 예: 데이터 입력 → 초안 30분 + 윤문 1시간 = 1.5시간."),
        ("시간 절감 추정 (%)", "예: 60%"),
        ("흔히 마주친 문제 · 해결", "예: 환각 → 사실 검증 1단계 추가."),
        ("다른 직군 적용 가능성", "예: 영업·마케팅에서도 사내 어휘집만 교체하면 동일하게 적용 가능."),
    ]

    for label, hint in fields:
        p = doc.add_paragraph()
        r = p.add_run(f"• {label}")
        r.bold = True
        h = doc.add_paragraph()
        hr = h.add_run(f"   {hint}")
        hr.font.size = Pt(10)
        hr.font.italic = True
        hr.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
        doc.add_paragraph("   [작성]")

    doc.add_heading("3. 빈 케이스 양식 (그대로 채워서 제출)", level=1)
    for i in range(1, 4):
        doc.add_heading(f"케이스 #{i}", level=2)
        for label, _ in fields:
            p = doc.add_paragraph()
            r = p.add_run(f"• {label}: ")
            r.bold = True
            p.add_run("[                                                                 ]")
        doc.add_paragraph()

    doc.add_heading("4. 분기 1회 컨테스트 안내", level=1)
    doc.add_paragraph(
        "분기말 우수 케이스 3건을 사내 발표 시간에 공유합니다. 우수 케이스는 사내 케이스 카드 갤러리(exec-4)에 등재되며, "
        "AX 챔피언 명단에 1년간 표기됩니다."
    )

    _add_page_number_footer(doc, FOOTER)
    out = os.path.join(OUTDIR, "case-collection-form.docx")
    doc.save(out)
    print(f"  wrote case-collection-form.docx")


def make_internal_guide_docx():
    doc = Document()
    _docx_set_margins(doc)
    _docx_title(
        doc,
        "[회사명] AI 사용 가이드 v1.0",
        "작성: [날짜]    ·    작성자: [AX 챔피언 이름]    ·    차기 갱신: [3개월 뒤]",
    )

    doc.add_heading("1. 우리 회사의 AX 비전", level=1)
    doc.add_paragraph(
        "[한 문장. 예: \"고객의 신뢰를 지키면서 제품 개발·고객 지원·콘텐츠 제작을 AI로 재설계해서, "
        "3년 뒤 1인당 매출 3배의 조직이 된다.\"]"
    )

    doc.add_heading("2. 사용 가능한 AI 도구", level=1)
    tbl = doc.add_table(rows=5, cols=5)
    tbl.style = "Light Grid Accent 1"
    hdr = tbl.rows[0].cells
    for i, h in enumerate(["도구", "등급", "사용 가능 데이터", "비용", "라이선스"]):
        hdr[i].text = h
    _bold_header_row(hdr)
    rows = [
        ("Claude (외부)", "메인", "공개·내부", "월 $20", "개인 결제"),
        ("ChatGPT (외부)", "보조", "공개·내부", "월 $20", "개인 결제"),
        ("Gemini Workspace", "메인 (Google 환경)", "공개·내부·기밀", "회사 라이선스", "회사 발급"),
        ("사내 LLM", "기밀 데이터용", "모든 등급", "인프라 운영비", "사내 운영"),
    ]
    for i, row in enumerate(rows, start=1):
        for j, v in enumerate(row):
            tbl.rows[i].cells[j].text = v

    doc.add_heading("3. 데이터 등급별 사용 규칙", level=1)
    for line in [
        "Lv1 공개: 모든 도구 OK",
        "Lv2 내부: 외부 LLM은 가공 후 OK",
        "Lv3 기밀: 사내 LLM만 OK",
        "Lv4 극비: 사내 LLM + 임원 승인",
    ]:
        doc.add_paragraph(line, style="List Bullet")

    doc.add_heading("4. 금지 행위", level=1)
    for line in [
        "극비 데이터(고객 개인정보·계약·재무 디테일)를 외부 LLM에 입력",
        "AI 출력물을 검증 없이 외부 송출 (외부 송출 게이트 통과 필수)",
        "정책 미승인 무료 도구 사용",
    ]:
        doc.add_paragraph(line, style="List Number")

    doc.add_heading("5. 외부 송출 게이트 (RGSB)", level=1)
    doc.add_paragraph(
        "외부에 나가는 모든 AI 협업 산출물은 다음 4관점 통과 후 송출합니다."
    )
    for line in [
        "R (Reason): 논리·근거 OK?",
        "G (Domain): 도메인 디테일 정확?",
        "S (Empathy): 독자 감정 흐름 OK?",
        "B (Reader): 최종 독자 첫 반응 OK?",
    ]:
        doc.add_paragraph(line, style="List Bullet")

    doc.add_heading("6. 사고 발생 시", level=1)
    doc.add_paragraph(
        "(1) 즉시 보고 → (2) 영향 평가 → (3) 외부 통보 (필요 시) → (4) 사후 분석 → (5) 재발 방지 룰 추가"
    )

    doc.add_heading("7. 챔피언 명단", level=1)
    doc.add_paragraph("[부서] [이름] — 슬랙 @[handle]")

    doc.add_heading("8. 분기 회고", level=1)
    for line in [
        "분기말 1시간 워크숍",
        "새 케이스 5개 공유",
        "가이드 v1.1 / v1.2 ... 업데이트",
    ]:
        doc.add_paragraph(line, style="List Bullet")

    _add_page_number_footer(doc, FOOTER)
    out = os.path.join(OUTDIR, "internal-guide-template.docx")
    doc.save(out)
    print(f"  wrote internal-guide-template.docx")


ONBOARDING_JOBS = [
    {
        "name": "경영기획·전략",
        "actions": ["AX 비전 문서 1회 통독", "사내 어휘집·과거 보고서 Claude Project에 업로드", "주간 보고서 1건을 Claude로 초안 → 사람 윤문"],
        "workflows": ["임원 보고서 자동 초안", "전략 메모 5-Color Harness 검토", "이사회 자료 RGSB 게이트 통과"],
        "tools": "Claude (메인) + 사내 어휘집 + 5-Color Harness 패턴",
        "faqs": [
            ("Q. 임원이 \"AI로 썼냐\"고 물으면?", "A. \"초안은 AI, 윤문·판단은 제가 했습니다.\"라고 답하면 됩니다."),
            ("Q. 컨피던셜 자료는?", "A. Lv3-4는 사내 LLM 또는 Workspace Gemini만 사용."),
            ("Q. 평가는?", "A. 시간 절감 + 9.5점 품질 기준 유지가 핵심 KPI."),
        ],
    },
    {
        "name": "영업",
        "actions": ["고객 메일 톤 1주일 ChatGPT 시도", "RFP 응답서 Claude Project로 초안 1건", "고객 리서치 Perplexity로 미팅 30분 전 5분"],
        "workflows": ["고객 메일 톤 맞춤 작성", "RFP 응답서 자동 초안", "고객 미팅 사전 리서치"],
        "tools": "ChatGPT (메일) + Claude (RFP) + Perplexity (리서치)",
        "faqs": [
            ("Q. 고객사 이름 외부 LLM에 입력해도 되나?", "A. 일반 정보(공개 사이트 기반)는 OK, 비공개 영업 정보는 사내 LLM."),
            ("Q. 메일 톤이 이상하면?", "A. 사내 어휘집을 시스템 프롬프트에 한 번 박으면 일관성 확보."),
            ("Q. 견적·계약 검토는?", "A. 외부 LLM 금지. 사내 LLM 또는 법무팀 확인."),
        ],
    },
    {
        "name": "마케팅",
        "actions": ["브랜드 가이드 ChatGPT Project에 업로드", "SNS 카피 5종 동시 생성 시도", "블로그 SEO 키워드 → Claude 초안"],
        "workflows": ["SNS 카피 5종 동시 생성", "이메일 캠페인 A/B 안 작성", "블로그 SEO 초안 작성"],
        "tools": "ChatGPT (카피) + Claude (장문) + Midjourney/Adobe (이미지)",
        "faqs": [
            ("Q. 이미지 생성 도구 선택은?", "A. 브랜드 톤 일관성 위해 Adobe Firefly 또는 사내 Stable Diffusion."),
            ("Q. 외부 송출 전 점검?", "A. 모든 외부 카피는 RGSB 게이트 통과."),
            ("Q. 톤이 광고스럽다?", "A. 사내 어휘집 + 과거 우수 카피 3개를 few-shot 예시로."),
        ],
    },
    {
        "name": "HR",
        "actions": ["JD 템플릿 Claude Project에 업로드", "면접 질문 5세트 30분 내 생성", "온보딩 안내문 1인 맞춤 시도"],
        "workflows": ["JD 채용 공고 자동 작성", "면접 질문 셋 자동 생성", "온보딩 안내문 개인 맞춤"],
        "tools": "Claude (장문 문서) + ChatGPT (개인 맞춤)",
        "faqs": [
            ("Q. 후보자 이력서를 외부 LLM에 입력?", "A. 금지. 사내 LLM 또는 가명화 후만."),
            ("Q. 차별적 표현 검증은?", "A. 출력 후 \"이 JD에서 차별적 표현 가능성 3가지\"를 항상 추가 질문."),
            ("Q. 평가 인터뷰 결과 분석?", "A. 사내 LLM만. 외부 LLM 금지."),
        ],
    },
    {
        "name": "재무·회계",
        "actions": ["월 결산 데이터 → Claude 변동 요약", "예산 차이 설명문 30분 내 작성", "투자 검토 메모 1차 초안 시도"],
        "workflows": ["월 결산 변동 분석 요약", "예산 차이 설명문 자동 작성", "투자 검토 메모 1차 초안"],
        "tools": "Claude (장문) + ChatGPT (요약) + 사내 LLM (실데이터)",
        "faqs": [
            ("Q. 실재무 데이터를 외부 LLM에?", "A. 절대 금지. 가명화/요약화 후 또는 사내 LLM."),
            ("Q. 숫자 검증은?", "A. 모든 숫자는 원본 엑셀과 1대1 매칭 확인 필수."),
            ("Q. 외부 감사용 메모는?", "A. 외부 송출 게이트 + 임원 승인."),
        ],
    },
    {
        "name": "개발",
        "actions": ["CLAUDE.md를 프로젝트 루트에 추가", "Cursor 또는 Claude Code 설치", "PR 1건을 Claude Code 리뷰로 시도"],
        "workflows": ["PR 리뷰 자동 보조", "테스트 코드 자동 생성", "버그 리포트 → 재현 코드"],
        "tools": "Claude Code + Cursor + GitHub Copilot (보조)",
        "faqs": [
            ("Q. 프로덕션 DB 쿼리?", "A. CLAUDE.md에 절대 금지로 명시. 읽기 전용 복제본만."),
            ("Q. 시크릿 관리?", "A. 환경 변수 + 사내 시크릿 매니저. AI에 직접 입력 금지."),
            ("Q. 외부 LLM에 사내 코드?", "A. 비공개 비즈니스 로직은 사내 LLM. 일반 라이브러리는 OK."),
        ],
    },
    {
        "name": "디자인",
        "actions": ["브랜드 톤 ChatGPT Project에 업로드", "시안 텍스트 자동 채우기 1건 시도", "사용자 인터뷰 녹취 → Claude 인사이트 추출"],
        "workflows": ["디자인 시안 텍스트 자동 채우기", "사용자 인터뷰 인사이트 추출", "디자인 리뷰 코멘트 정리"],
        "tools": "ChatGPT (텍스트) + Claude (분석) + Figma AI (시안)",
        "faqs": [
            ("Q. 사용자 인터뷰 녹취 외부 LLM에?", "A. 가명화 후 OK. 회사명·개인정보 제거 필수."),
            ("Q. 디자인 톤 일관성?", "A. 브랜드 가이드 + 과거 우수 시안 few-shot 예시로."),
            ("Q. AI 이미지 라이선스?", "A. Adobe Firefly = 상업 사용 OK. 다른 도구는 사내 법무 확인."),
        ],
    },
    {
        "name": "고객지원",
        "actions": ["FAQ DB를 사내 LLM RAG에 업로드", "1차 자동 답변 1주일 운영", "주간 불만 톤 분석 30분 시도"],
        "workflows": ["FAQ 답변 1차 자동화", "고객 불만 톤 분석·우선순위", "복잡 케이스 대응 시나리오 생성"],
        "tools": "사내 LLM (1차 자동) + Claude (분석) + ChatGPT (다국어)",
        "faqs": [
            ("Q. 고객 개인정보?", "A. 사내 LLM만. 외부 LLM에 입력 절대 금지."),
            ("Q. 1차 자동 답변 검증?", "A. 모든 자동 답변은 사람 검토 후 전송 (RGSB 게이트)."),
            ("Q. 다국어는?", "A. Claude·ChatGPT 모두 강력. 단 전문 용어는 사내 어휘집 첨부."),
        ],
    },
]


def make_onboarding_1pagers():
    doc = Document()
    _docx_set_margins(doc, inches=0.8)
    _docx_title(
        doc,
        "직군별 AX 온보딩 1-pager 8종",
        "AX Guide for Enterprise · 첫 1주일을 위한 통합 자료",
    )

    for idx, job in enumerate(ONBOARDING_JOBS):
        if idx > 0:
            doc.add_paragraph().add_run().add_break(WD_BREAK.PAGE)

        h = doc.add_heading(f"[{job['name']}] AX 온보딩 1-pager", level=1)

        doc.add_heading("첫 주 액션 3개", level=2)
        for a in job["actions"]:
            doc.add_paragraph(f"☐ {a}")

        doc.add_heading("핵심 워크플로우 3개 (exec-1과 매칭)", level=2)
        for w in job["workflows"]:
            doc.add_paragraph(w, style="List Number")

        doc.add_heading("추천 도구", level=2)
        doc.add_paragraph(job["tools"])

        doc.add_heading("자주 묻는 질문 3개", level=2)
        for q, a in job["faqs"]:
            p = doc.add_paragraph()
            r = p.add_run(q)
            r.bold = True
            doc.add_paragraph(a)

        doc.add_heading("챔피언 연락처", level=2)
        doc.add_paragraph("이름: [                ]   ·   슬랙: @[                ]   ·   메일: [                ]")

    _add_page_number_footer(doc, FOOTER)
    out = os.path.join(OUTDIR, "onboarding-1pagers.docx")
    doc.save(out)
    print(f"  wrote onboarding-1pagers.docx")


def make_monthly_retrospective():
    doc = Document()
    _docx_set_margins(doc)
    _docx_title(
        doc,
        "분기 회고 워크숍 양식",
        "분기 1회 · 1시간 · AX 챔피언 + 부서 대표",
    )

    doc.add_heading("1. 목적", level=1)
    doc.add_paragraph(
        "분기 1회 1시간 워크숍에서 KPI를 점검하고, 우수 케이스를 공유하며, 다음 분기 우선순위를 정리합니다. "
        "회고는 AX 가이드 갱신(v1.1, v1.2...)의 단일 근거가 됩니다."
    )

    doc.add_heading("2. 진행 순서 (총 60분)", level=1)
    schedule = [
        ("5분", "도입", "지난 분기 핵심 메시지 1줄 요약"),
        ("20분", "KPI 리뷰", "사용자 수 / 시간 절감 / 사고 / 가이드 위반"),
        ("20분", "우수 케이스 공유", "3-5개 케이스, 1건당 3-4분"),
        ("10분", "다음 분기 우선순위", "TOP 3 액션 결정"),
        ("5분", "정리", "기록자 발표 + 다음 회고 날짜 확정"),
    ]
    tbl = doc.add_table(rows=1, cols=3)
    tbl.style = "Light Grid Accent 1"
    hdr = tbl.rows[0].cells
    hdr[0].text = "시간"
    hdr[1].text = "단계"
    hdr[2].text = "내용"
    _bold_header_row(hdr)
    for t, stage, content in schedule:
        row = tbl.add_row().cells
        row[0].text = t
        row[1].text = stage
        row[2].text = content

    doc.add_heading("3. 체크리스트 양식", level=1)
    for label in [
        "이번 분기 활성 사용자 수: [          ]",
        "이번 분기 평균 시간 절감 (%): [          ]",
        "사고 발생 건수: [    ] / 외부 노출 건수: [    ]",
        "가이드 위반 건수: [    ] (1차 경고 [    ] / 2차 [    ])",
        "신규 등록 케이스 수: [    ]",
        "분기말 챔피언 명단 갱신 필요 여부: [예/아니오]",
        "다음 분기 도입 도구·라이선스 변경: [          ]",
    ]:
        doc.add_paragraph(f"☐ {label}")

    doc.add_heading("4. 다음 분기 액션 5개", level=1)
    for i in range(1, 6):
        p = doc.add_paragraph()
        r = p.add_run(f"액션 #{i}: ")
        r.bold = True
        p.add_run("[                                                  ] / 담당: [        ] / 시한: [        ]")

    _add_page_number_footer(doc, FOOTER)
    out = os.path.join(OUTDIR, "monthly-retrospective.docx")
    doc.save(out)
    print(f"  wrote monthly-retrospective.docx")


def make_tool_decision_tree():
    doc = Document()
    _docx_set_margins(doc)
    _docx_title(
        doc,
        "AI 도구 결정 트리",
        "5단계 질문으로 우리 회사에 맞는 도구 1개를 고른다",
    )

    doc.add_heading("결정 트리", level=1)
    doc.add_paragraph(
        "아래 질문을 순서대로 따라가며 답하면, 직군과 업무 컨텍스트에 맞는 AI 도구를 한 가지로 좁힐 수 있습니다."
    )

    steps = [
        ("Q1.", "긴 문서·보고서·장문 분석이 주 업무?",
         "→ YES: Claude (장문·문서 강점)",
         "→ NO: Q2로"),
        ("Q2.", "Google Workspace (Docs·Sheets·Gmail) 환경?",
         "→ YES: Gemini Workspace (Drive 직접 연동)",
         "→ NO: Q3로"),
        ("Q3.", "Microsoft 365 (Word·Excel·Outlook) 환경?",
         "→ YES: Microsoft Copilot (M365 통합)",
         "→ NO: Q4로"),
        ("Q4.", "사내 데이터 Lv3-4 (기밀·극비) 다룸?",
         "→ YES: 사내 LLM (데이터 외부 유출 차단)",
         "→ NO: Q5로"),
        ("Q5.", "코드 작업이 주 업무?",
         "→ YES: Claude Code 또는 Cursor",
         "→ NO: ChatGPT 또는 Perplexity (범용·리서치)"),
    ]

    for q, prompt, yes, no in steps:
        h = doc.add_paragraph()
        hr = h.add_run(f"{q}  {prompt}")
        hr.bold = True
        hr.font.size = Pt(13)

        y = doc.add_paragraph(f"    {yes}")
        for r in y.runs:
            r.font.size = Pt(11)

        n = doc.add_paragraph(f"    {no}")
        for r in n.runs:
            r.font.size = Pt(11)
        doc.add_paragraph()

    doc.add_heading("주의", level=1)
    for line in [
        "도구 선택은 한 번이 아니라 분기마다 재검토.",
        "사용자 라이선스·데이터 정책은 분기 1회 확인.",
        "코드·기밀 데이터는 외부 LLM 금지가 기본.",
        "고객사 결정 등 외부 송출 자산은 RGSB 게이트 통과 후 송출.",
    ]:
        doc.add_paragraph(line, style="List Bullet")

    _add_page_number_footer(doc, FOOTER)
    out = os.path.join(OUTDIR, "tool-decision-tree.docx")
    doc.save(out)
    print(f"  wrote tool-decision-tree.docx")


# ===================== PPTX =====================

NAVY = PColor(0x11, 0x18, 0x27)
GREY = PColor(0x6B, 0x72, 0x80)
LIGHT_GREY = PColor(0xE5, 0xE7, 0xEB)
BG_LIGHT = PColor(0xF9, 0xFA, 0xFB)
ACCENT = PColor(0xF5, 0x9E, 0x0B)
GREEN = PColor(0x10, 0xB9, 0x81)


def _slide_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def _add_text(slide, left, top, width, height, text, *, size=14, bold=False, color=NAVY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Pin(left), Pin(top), Pin(width), Pin(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = PPt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Inter"
    return tb


def _add_footer(slide, prs):
    _add_text(slide, 0.3, prs.slide_height.inches - 0.4, prs.slide_width.inches - 0.6, 0.3,
              f"[회사명]    ·    {FOOTER}",
              size=9, color=GREY, align=PP_ALIGN.CENTER)


def _box(slide, left, top, width, height, fill=BG_LIGHT, line=LIGHT_GREY):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pin(left), Pin(top), Pin(width), Pin(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = PPt(0.5)
    shp.shadow.inherit = False
    return shp


def make_case_presentation():
    prs = Presentation()
    prs.slide_width = Pin(13.333)
    prs.slide_height = Pin(7.5)

    # Slide 1 — Title
    s = _slide_blank(prs)
    _add_text(s, 0.7, 2.5, 12, 1.2,
              "AI 협업 사례: [케이스 제목]",
              size=42, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    _add_text(s, 0.7, 3.8, 12, 0.6,
              "[작성자 / 부서 / 날짜]",
              size=18, color=GREY, align=PP_ALIGN.LEFT)
    _add_text(s, 0.7, 4.5, 12, 0.5,
              "AX Guide for Enterprise · 사내 케이스 발표 템플릿",
              size=14, color=GREY, align=PP_ALIGN.LEFT)
    _add_footer(s, prs)

    # Slide 2 — Before
    s = _slide_blank(prs)
    _add_text(s, 0.5, 0.3, 12, 0.6, "Before — 현재 일하는 방식",
              size=28, bold=True, color=NAVY)
    _add_text(s, 0.5, 1.0, 12, 0.4,
              "AI를 도입하기 전 우리가 어떻게 일하고 있었는가",
              size=13, color=GREY)

    _box(s, 0.5, 1.7, 6, 4.5)
    _add_text(s, 0.8, 1.9, 5.5, 0.5, "기존 프로세스", size=16, bold=True, color=NAVY)
    _add_text(s, 0.8, 2.5, 5.5, 3.5,
              "• 1단계: [구체적 행동 + 시간]\n"
              "• 2단계: [구체적 행동 + 시간]\n"
              "• 3단계: [구체적 행동 + 시간]\n"
              "• 4단계: [구체적 행동 + 시간]\n\n"
              "총 소요 시간: [X시간 / X일]",
              size=14, color=NAVY)

    _box(s, 7.0, 1.7, 5.8, 4.5, fill=PColor(0xFE, 0xE2, 0xE2), line=PColor(0xFC, 0xA5, 0xA5))
    _add_text(s, 7.3, 1.9, 5.3, 0.5, "병목 · 문제점", size=16, bold=True, color=NAVY)
    _add_text(s, 7.3, 2.5, 5.3, 3.5,
              "• 시간이 가장 많이 드는 단계: [    ]\n"
              "• 품질이 흔들리는 지점: [    ]\n"
              "• 사람 의존도가 큰 부분: [    ]\n"
              "• 반복 작업의 비중: [    ]",
              size=14, color=NAVY)

    _add_footer(s, prs)

    # Slide 3 — 도구·프롬프트·검증
    s = _slide_blank(prs)
    _add_text(s, 0.5, 0.3, 12, 0.6, "도구 · 프롬프트 · 검증",
              size=28, bold=True, color=NAVY)
    _add_text(s, 0.5, 1.0, 12, 0.4,
              "어떤 도구·어떤 프롬프트 구조·어떤 검증 패턴을 썼는가",
              size=13, color=GREY)

    _box(s, 0.5, 1.7, 4.0, 4.5)
    _add_text(s, 0.7, 1.9, 3.6, 0.5, "도구", size=16, bold=True, color=NAVY)
    _add_text(s, 0.7, 2.5, 3.6, 3.5,
              "• 메인: [Claude / GPT / Gemini]\n"
              "• 보조: [Perplexity 등]\n"
              "• 라이선스: [개인/팀]\n"
              "• 데이터 등급: [Lv1-4]",
              size=14, color=NAVY)

    _box(s, 4.7, 1.7, 4.0, 4.5)
    _add_text(s, 4.9, 1.9, 3.6, 0.5, "프롬프트 구조", size=16, bold=True, color=NAVY)
    _add_text(s, 4.9, 2.5, 3.6, 3.5,
              "• 시스템 프롬프트: [있음/없음]\n"
              "• 사내 어휘집: [첨부 여부]\n"
              "• 단계 분리: [몇 단계]\n"
              "• Few-shot 예시: [있음/없음]",
              size=14, color=NAVY)

    _box(s, 8.9, 1.7, 3.9, 4.5, fill=PColor(0xD1, 0xFA, 0xE5), line=GREEN)
    _add_text(s, 9.1, 1.9, 3.5, 0.5, "검증 패턴", size=16, bold=True, color=NAVY)
    _add_text(s, 9.1, 2.5, 3.5, 3.5,
              "• 자기검증 루프 회수\n"
              "• 5-Color Harness 사용?\n"
              "• RGSB 게이트 통과?\n"
              "• 사람 검수 시간: [    ]",
              size=14, color=NAVY)

    _add_footer(s, prs)

    # Slide 4 — After
    s = _slide_blank(prs)
    _add_text(s, 0.5, 0.3, 12, 0.6, "After — 바뀐 일하는 방식",
              size=28, bold=True, color=NAVY)
    _add_text(s, 0.5, 1.0, 12, 0.4,
              "AI 도입 후 우리가 어떻게 일하게 되었는가",
              size=13, color=GREY)

    _box(s, 0.5, 1.7, 6, 4.5)
    _add_text(s, 0.8, 1.9, 5.5, 0.5, "변화된 프로세스", size=16, bold=True, color=NAVY)
    _add_text(s, 0.8, 2.5, 5.5, 3.5,
              "• 1단계: [AI가 수행 + 시간]\n"
              "• 2단계: [사람이 수행 + 시간]\n"
              "• 3단계: [AI 검증]\n"
              "• 4단계: [사람 최종 승인]\n\n"
              "총 소요 시간: [X시간 / X일]",
              size=14, color=NAVY)

    _box(s, 7.0, 1.7, 5.8, 4.5, fill=PColor(0xD1, 0xFA, 0xE5), line=GREEN)
    _add_text(s, 7.3, 1.9, 5.3, 0.5, "품질·문화 변화", size=16, bold=True, color=NAVY)
    _add_text(s, 7.3, 2.5, 5.3, 3.5,
              "• 품질: [9.5점 기준 만족도]\n"
              "• 사람의 역할 변화: [수행자 → 검수자]\n"
              "• 다른 일에 쓸 수 있게 된 시간\n"
              "• 사고·오류 발생 여부",
              size=14, color=NAVY)

    _add_footer(s, prs)

    # Slide 5 — KPI + 확산
    s = _slide_blank(prs)
    _add_text(s, 0.5, 0.3, 12, 0.6, "시간 절감 + 다음 적용",
              size=28, bold=True, color=NAVY)
    _add_text(s, 0.5, 1.0, 12, 0.4,
              "정량 효과와 어디로 확산할지",
              size=13, color=GREY)

    # 3 KPI boxes
    kpis = [
        ("시간 절감", "[ 60% ]", "Before X시간 → After Y시간"),
        ("품질 등급", "[ 9.5 / 10 ]", "5-Color 평균"),
        ("월간 절감 인시", "[ X 인시 ]", "팀 단위 환산"),
    ]
    for i, (label, value, sub) in enumerate(kpis):
        left = 0.5 + i * 4.3
        _box(s, left, 1.7, 4.0, 2.2, fill=PColor(0xFE, 0xF3, 0xC7), line=ACCENT)
        _add_text(s, left + 0.2, 1.85, 3.6, 0.5, label, size=14, bold=True, color=NAVY)
        _add_text(s, left + 0.2, 2.4, 3.6, 0.9, value, size=28, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        _add_text(s, left + 0.2, 3.3, 3.6, 0.5, sub, size=11, color=GREY, align=PP_ALIGN.CENTER)

    _box(s, 0.5, 4.2, 12.3, 2.4)
    _add_text(s, 0.7, 4.3, 11.9, 0.5, "확산 계획", size=16, bold=True, color=NAVY)
    _add_text(s, 0.7, 4.9, 11.9, 1.6,
              "• 1순위 적용 대상: [부서·역할]\n"
              "• 필요한 사전 작업: [어휘집 / 라이선스 / 교육]\n"
              "• 예상 차이: [원안과 다른 점 + 추가 검증 필요한 부분]\n"
              "• 다음 발표·등록 일정: [날짜]",
              size=14, color=NAVY)

    _add_footer(s, prs)

    out = os.path.join(OUTDIR, "case-presentation.pptx")
    prs.save(out)
    print(f"  wrote case-presentation.pptx")


# ===================== MAIN =====================

def main():
    print(f"Building P3 assets in {OUTDIR} ...")
    print("\n[Markdown]")
    make_markdown_files()

    print("\n[Excel]")
    make_tool_pricing_comparison()
    make_cases_30_list()

    print("\n[Word]")
    make_case_collection_form()
    make_internal_guide_docx()
    make_onboarding_1pagers()
    make_monthly_retrospective()
    make_tool_decision_tree()

    print("\n[PowerPoint]")
    make_case_presentation()

    print("\nDone.")


if __name__ == "__main__":
    main()
