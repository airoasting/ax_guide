"""
Build P2 (Strategy track) download assets for AX Guide for Enterprise.

Outputs (in same directory as this script):
- roi-calculator.xlsx
- impact-feasibility-matrix.pptx
- payback-model.xlsx
- ai-usage-policy.docx
- data-classification.xlsx

Run from repo root:
    python3 assets/downloads/_build_assets.py
"""

import os

from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from pptx import Presentation
from pptx.util import Inches as Pin, Pt as PPt, Emu
from pptx.dml.color import RGBColor as PColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTDIR = SCRIPT_DIR
FOOTER = "AX Guide for Enterprise · v0.2 / jade@linercorp.com"


# ---------- common style helpers ----------

THIN = Side(border_style="thin", color="CCCCCC")
BORDER = Border(left=THIN, right=THIN, top=THIN, bottom=THIN)

HEADER_FILL = PatternFill("solid", fgColor="1F2937")
SUBHEADER_FILL = PatternFill("solid", fgColor="E5E7EB")
ACCENT_FILL = PatternFill("solid", fgColor="FEF3C7")

HEADER_FONT = Font(name="Inter", size=12, bold=True, color="FFFFFF")
SUBHEADER_FONT = Font(name="Inter", size=11, bold=True, color="111827")
BODY_FONT = Font(name="Inter", size=11, color="111827")
HINT_FONT = Font(name="Inter", size=10, italic=True, color="6B7280")
FOOTER_FONT = Font(name="Inter", size=9, italic=True, color="6B7280")


def set_col_widths(ws, widths):
    for i, w in enumerate(widths, start=1):
        ws.column_dimensions[get_column_letter(i)].width = w


def put(ws, cell, value, font=None, fill=None, border=None, align=None):
    c = ws[cell]
    c.value = value
    if font:
        c.font = font
    if fill:
        c.fill = fill
    if border:
        c.border = border
    if align:
        c.alignment = align


# ---------- 1) ROI calculator ----------

def make_roi_calculator():
    wb = Workbook()

    # Sheet 1: ROI 계산
    ws = wb.active
    ws.title = "ROI 계산"
    set_col_widths(ws, [28, 22, 40])

    put(ws, "A1", "AX ROI 12개월 회수 모델", Font(name="Inter", size=16, bold=True, color="111827"))
    ws.merge_cells("A1:C1")
    ws.row_dimensions[1].height = 28

    put(ws, "A2", "회사명", SUBHEADER_FONT, SUBHEADER_FILL, BORDER)
    put(ws, "B2", "", BODY_FONT, None, BORDER)
    put(ws, "C2", "예: ABC주식회사", HINT_FONT, None, BORDER)

    put(ws, "A3", "작성일", SUBHEADER_FONT, SUBHEADER_FILL, BORDER)
    put(ws, "B3", "", BODY_FONT, None, BORDER)
    put(ws, "C3", "YYYY-MM-DD", HINT_FONT, None, BORDER)

    # Header row
    headers = ["항목", "값", "비고"]
    for i, h in enumerate(headers):
        put(ws, f"{get_column_letter(i+1)}5", h, HEADER_FONT, HEADER_FILL, BORDER,
            Alignment(horizontal="center", vertical="center"))
    ws.row_dimensions[5].height = 22

    center = Alignment(horizontal="center", vertical="center")
    left = Alignment(horizontal="left", vertical="center", wrap_text=True)

    rows = [
        ("대상 직군", "", "예: 마케팅팀"),
        ("인원 수", "", "명"),
        ("평균 시급 (만원)", "", "만원/시간"),
        ("주당 절감 시간 (1인)", "", "시간"),
        ("연 절감 시간 (1인)", "=B9*48", "시간 (=주당*48주)"),
        ("연 절감액 (1인, 만원)", "=B10*B8", "=연 절감 시간 * 시급"),
        ("연 절감액 (전체, 만원)", "=B11*B7", "=1인 절감액 * 인원수"),
        ("라이선스 비용 (연, 만원)", "", "AI 도구 라이선스 합계"),
        ("교육 비용 (1회, 만원)", "", "초기 교육 비용"),
        ("관리감독 비용 (연, 만원)", "", "거버넌스/모니터링"),
        ("총 비용 (만원)", "=B13+B14+B15", "라이선스+교육+관리"),
        ("순 ROI (만원)", "=B12-B16", "=전체 절감액 - 총 비용"),
        ("회수 기간 (개월)", '=IF(B12=0,"-",12*B16/B12)', "=12*비용/연절감액"),
    ]

    start_row = 6
    for i, (item, val, memo) in enumerate(rows):
        r = start_row + i
        put(ws, f"A{r}", item, BODY_FONT, None, BORDER, left)
        # highlight calculated rows
        is_calc = isinstance(val, str) and val.startswith("=")
        fill = ACCENT_FILL if is_calc else None
        put(ws, f"B{r}", val, BODY_FONT, fill, BORDER, center)
        put(ws, f"C{r}", memo, HINT_FONT, None, BORDER, left)

    # Note
    note_row = start_row + len(rows) + 2
    put(ws, f"A{note_row}",
        "노랑 셀은 자동 계산됩니다. 흰색 셀에 값을 채우면 ROI와 회수 기간이 자동 산출됩니다.",
        HINT_FONT)
    ws.merge_cells(f"A{note_row}:C{note_row}")

    footer_row = note_row + 2
    put(ws, f"A{footer_row}", FOOTER, FOOTER_FONT)
    ws.merge_cells(f"A{footer_row}:C{footer_row}")

    # Sheet 2: 직군별 추정 시간
    ws2 = wb.create_sheet("직군별 추정 시간")
    set_col_widths(ws2, [24, 24, 40])

    put(ws2, "A1", "직군별 주당 절감 시간 추정 (참고)",
        Font(name="Inter", size=14, bold=True, color="111827"))
    ws2.merge_cells("A1:C1")
    ws2.row_dimensions[1].height = 24

    headers2 = ["직군", "주당 절감 시간 (h)", "비고"]
    for i, h in enumerate(headers2):
        put(ws2, f"{get_column_letter(i+1)}3", h, HEADER_FONT, HEADER_FILL, BORDER, center)
    ws2.row_dimensions[3].height = 22

    jobs = [
        ("경영기획·전략", "", "보고서/분석 자동화 (권장 측정 범위 4-12h)"),
        ("영업", "", "제안서/이메일/CRM 입력 자동화"),
        ("마케팅", "", "콘텐츠 초안/SEO/소재 변형"),
        ("HR", "", "JD 작성/스크리닝/온보딩"),
        ("재무·회계", "", "리포트/분개 검토/비교 분석"),
        ("개발", "", "코드 생성/리뷰/문서화"),
        ("디자인", "", "시안 변형/카피/리서치"),
        ("고객지원", "", "응답 초안/FAQ/요약"),
    ]
    for i, (job, hours, memo) in enumerate(jobs):
        r = 4 + i
        put(ws2, f"A{r}", job, BODY_FONT, None, BORDER, left)
        put(ws2, f"B{r}", hours, BODY_FONT, None, BORDER, center)
        put(ws2, f"C{r}", memo, HINT_FONT, None, BORDER, left)

    note2 = 4 + len(jobs) + 2
    put(ws2, f"A{note2}",
        "추정 시간은 4-12h 범위에서 사내가 직접 측정한 값을 넣으세요. 추정만으로 ROI를 산출하면 신뢰도가 떨어집니다.",
        HINT_FONT)
    ws2.merge_cells(f"A{note2}:C{note2}")

    put(ws2, f"A{note2+2}", FOOTER, FOOTER_FONT)
    ws2.merge_cells(f"A{note2+2}:C{note2+2}")

    # Sheet 3: 안내
    ws3 = wb.create_sheet("안내")
    set_col_widths(ws3, [100])
    put(ws3, "A1", "사용 안내", Font(name="Inter", size=16, bold=True, color="111827"))
    guide = (
        "이 시트는 9.5점 ROI 계산 모델입니다.\n\n"
        "정확한 측정 없이 추정만 채워 넣지 마세요. "
        "첫 분기에는 1개 직군 파일럿 측정 후 다른 직군으로 확장하는 것을 권장합니다.\n\n"
        "권장 사용 흐름:\n"
        "  1) 파일럿 직군 1개 선정 (예: 마케팅팀 5명)\n"
        "  2) 4주간 도입 전후 시간 측정 (베이스라인 vs 도입 후)\n"
        "  3) '직군별 추정 시간' 시트에 실측값 기록\n"
        "  4) 'ROI 계산' 시트에 값 입력 → 회수 기간 산출\n"
        "  5) 다른 직군으로 확장 시 동일 절차 반복\n\n"
        "주의:\n"
        "  - 라이선스 비용은 연간 총액으로 환산하세요.\n"
        "  - 관리감독 비용에는 거버넌스 운영, 품질 모니터링, 사고 대응이 포함됩니다.\n"
        "  - 회수 기간이 18개월을 넘으면 도입 우선순위를 재검토하세요.\n\n"
        f"— {FOOTER}"
    )
    put(ws3, "A3", guide, Font(name="Inter", size=11, color="111827"),
        align=Alignment(wrap_text=True, vertical="top"))
    ws3.row_dimensions[3].height = 320

    out = os.path.join(OUTDIR, "roi-calculator.xlsx")
    wb.save(out)
    print(f"  wrote {out}")


# ---------- 2) Impact x Feasibility matrix (PPTX) ----------

def make_impact_matrix():
    prs = Presentation()
    prs.slide_width = Pin(13.333)
    prs.slide_height = Pin(7.5)

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Background white (default)

    # Title
    title_box = slide.shapes.add_textbox(Pin(0.5), Pin(0.3), Pin(12.3), Pin(0.7))
    tf = title_box.text_frame
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    p.text = "영향도 × 실현가능성 우선순위 매트릭스"
    p.font.size = PPt(26)
    p.font.bold = True
    p.font.color.rgb = PColor(0x11, 0x18, 0x27)
    p.alignment = PP_ALIGN.LEFT

    # Subtitle
    sub_box = slide.shapes.add_textbox(Pin(0.5), Pin(1.0), Pin(12.3), Pin(0.4))
    sp = sub_box.text_frame.paragraphs[0]
    sp.text = "AX 과제를 4분면에 배치해 의사결정 우선순위를 결정합니다."
    sp.font.size = PPt(12)
    sp.font.color.rgb = PColor(0x6B, 0x72, 0x80)

    # Matrix area
    # area: x=2.0, y=1.6, width=9.5, height=5.0
    ax = 2.0
    ay = 1.6
    aw = 9.5
    ah = 5.0
    cw = aw / 2
    ch = ah / 2

    # Quadrant definitions: (col, row, fill_rgb, title, subtitle)
    # row 0 = top (high impact), row 1 = bottom (low impact)
    # col 0 = left (low feasibility), col 1 = right (high feasibility)
    quads = [
        # 좌상: High Impact, Low Feasibility → 노랑
        (0, 0, (0xFE, 0xF3, 0xC7), "장기 R&D", "High Impact · Low Feasibility"),
        # 우상: High Impact, High Feasibility → 녹색
        (1, 0, (0xD1, 0xFA, 0xE5), "즉시 파일럿", "High Impact · High Feasibility"),
        # 좌하: Low Impact, Low Feasibility → 진한 회색
        (0, 1, (0x9C, 0xA3, 0xAF), "안 함", "Low Impact · Low Feasibility"),
        # 우하: Low Impact, High Feasibility → 연한 회색
        (1, 1, (0xE5, 0xE7, 0xEB), "소규모 도구 도입", "Low Impact · High Feasibility"),
    ]

    for col, row, rgb, title, sub in quads:
        x = ax + col * cw
        y = ay + row * ch
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pin(x), Pin(y), Pin(cw), Pin(ch))
        shp.fill.solid()
        shp.fill.fore_color.rgb = PColor(*rgb)
        shp.line.color.rgb = PColor(0xFF, 0xFF, 0xFF)
        shp.line.width = PPt(2)
        # remove shape default text behavior; add text via shape
        tf = shp.text_frame
        tf.margin_left = Pin(0.25)
        tf.margin_top = Pin(0.25)
        tf.margin_right = Pin(0.25)
        tf.margin_bottom = Pin(0.25)
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = PPt(22)
        p1.font.bold = True
        # Dark text on light fills, white text on dark fills
        text_color = (0xFF, 0xFF, 0xFF) if rgb == (0x9C, 0xA3, 0xAF) else (0x11, 0x18, 0x27)
        p1.font.color.rgb = PColor(*text_color)
        p1.alignment = PP_ALIGN.LEFT
        p2 = tf.add_paragraph()
        p2.text = sub
        p2.font.size = PPt(11)
        p2.font.color.rgb = PColor(*text_color)
        p2.alignment = PP_ALIGN.LEFT

    # Y-axis label (영향도 →, with Low/High markers)
    # Vertical text on left
    y_label = slide.shapes.add_textbox(Pin(0.3), Pin(ay), Pin(1.5), Pin(ah))
    ylt = y_label.text_frame
    ylt.word_wrap = True
    ylt.vertical_anchor = MSO_ANCHOR.MIDDLE
    yp = ylt.paragraphs[0]
    yp.text = "↑ 영향도 (Impact)"
    yp.font.size = PPt(13)
    yp.font.bold = True
    yp.font.color.rgb = PColor(0x37, 0x41, 0x51)
    yp.alignment = PP_ALIGN.CENTER

    # High marker (top)
    yhigh = slide.shapes.add_textbox(Pin(0.3), Pin(ay - 0.05), Pin(1.5), Pin(0.3))
    yhp = yhigh.text_frame.paragraphs[0]
    yhp.text = "High"
    yhp.font.size = PPt(10)
    yhp.font.color.rgb = PColor(0x6B, 0x72, 0x80)
    yhp.alignment = PP_ALIGN.CENTER

    # Low marker (bottom)
    ylow = slide.shapes.add_textbox(Pin(0.3), Pin(ay + ah - 0.3), Pin(1.5), Pin(0.3))
    ylp = ylow.text_frame.paragraphs[0]
    ylp.text = "Low"
    ylp.font.size = PPt(10)
    ylp.font.color.rgb = PColor(0x6B, 0x72, 0x80)
    ylp.alignment = PP_ALIGN.CENTER

    # X-axis label (실현가능성 →)
    x_label = slide.shapes.add_textbox(Pin(ax), Pin(ay + ah + 0.15), Pin(aw), Pin(0.4))
    xlt = x_label.text_frame
    xlp = xlt.paragraphs[0]
    xlp.text = "실현가능성 (Feasibility) →"
    xlp.font.size = PPt(13)
    xlp.font.bold = True
    xlp.font.color.rgb = PColor(0x37, 0x41, 0x51)
    xlp.alignment = PP_ALIGN.CENTER

    # X markers Low / High
    xlow = slide.shapes.add_textbox(Pin(ax - 0.05), Pin(ay + ah + 0.55), Pin(1.0), Pin(0.3))
    xlp2 = xlow.text_frame.paragraphs[0]
    xlp2.text = "Low"
    xlp2.font.size = PPt(10)
    xlp2.font.color.rgb = PColor(0x6B, 0x72, 0x80)

    xhigh = slide.shapes.add_textbox(Pin(ax + aw - 1.0 + 0.05), Pin(ay + ah + 0.55), Pin(1.0), Pin(0.3))
    xhp = xhigh.text_frame.paragraphs[0]
    xhp.text = "High"
    xhp.font.size = PPt(10)
    xhp.font.color.rgb = PColor(0x6B, 0x72, 0x80)
    xhp.alignment = PP_ALIGN.RIGHT

    # Footer
    footer_box = slide.shapes.add_textbox(Pin(0.5), Pin(7.1), Pin(12.3), Pin(0.3))
    fp = footer_box.text_frame.paragraphs[0]
    fp.text = FOOTER
    fp.font.size = PPt(9)
    fp.font.italic = True
    fp.font.color.rgb = PColor(0x6B, 0x72, 0x80)
    fp.alignment = PP_ALIGN.RIGHT

    out = os.path.join(OUTDIR, "impact-feasibility-matrix.pptx")
    prs.save(out)
    print(f"  wrote {out}")


# ---------- 3) Payback model ----------

def make_payback_model():
    wb = Workbook()
    ws = wb.active
    ws.title = "12개월 회수 모델"

    set_col_widths(ws, [22] + [11] * 12 + [14])

    put(ws, "A1", "12개월 누적 회수 모델",
        Font(name="Inter", size=16, bold=True, color="111827"))
    ws.merge_cells("A1:N1")
    ws.row_dimensions[1].height = 26

    # Meta rows
    put(ws, "A3", "회사명", SUBHEADER_FONT, SUBHEADER_FILL, BORDER)
    put(ws, "B3", "", BODY_FONT, None, BORDER)
    ws.merge_cells("B3:E3")
    put(ws, "F3", "작성일", SUBHEADER_FONT, SUBHEADER_FILL, BORDER)
    put(ws, "G3", "", BODY_FONT, None, BORDER)
    ws.merge_cells("G3:I3")
    put(ws, "J3", "대상 직군", SUBHEADER_FONT, SUBHEADER_FILL, BORDER)
    put(ws, "K3", "", BODY_FONT, None, BORDER)
    ws.merge_cells("K3:N3")

    # Month header row
    put(ws, "A5", "항목", HEADER_FONT, HEADER_FILL, BORDER,
        Alignment(horizontal="center", vertical="center"))
    for m in range(1, 13):
        col = get_column_letter(1 + m)
        put(ws, f"{col}5", f"M{m}", HEADER_FONT, HEADER_FILL, BORDER,
            Alignment(horizontal="center", vertical="center"))
    put(ws, "N5", "합계", HEADER_FONT, HEADER_FILL, BORDER,
        Alignment(horizontal="center", vertical="center"))
    ws.row_dimensions[5].height = 22

    center = Alignment(horizontal="center", vertical="center")
    left = Alignment(horizontal="left", vertical="center")

    # Row 6: 절감액
    put(ws, "A6", "절감액 (만원)", BODY_FONT, None, BORDER, left)
    for m in range(1, 13):
        col = get_column_letter(1 + m)
        put(ws, f"{col}6", "", BODY_FONT, None, BORDER, center)
    put(ws, "N6", "=SUM(B6:M6)", BODY_FONT, ACCENT_FILL, BORDER, center)

    # Row 7: 비용
    put(ws, "A7", "비용 (만원)", BODY_FONT, None, BORDER, left)
    for m in range(1, 13):
        col = get_column_letter(1 + m)
        put(ws, f"{col}7", "", BODY_FONT, None, BORDER, center)
    put(ws, "N7", "=SUM(B7:M7)", BODY_FONT, ACCENT_FILL, BORDER, center)

    # Row 8: 월 순이익 = 절감 - 비용
    put(ws, "A8", "월 순이익 (만원)", BODY_FONT, None, BORDER, left)
    for m in range(1, 13):
        col = get_column_letter(1 + m)
        put(ws, f"{col}8", f"={col}6-{col}7", BODY_FONT, ACCENT_FILL, BORDER, center)
    put(ws, "N8", "=SUM(B8:M8)", BODY_FONT, ACCENT_FILL, BORDER, center)

    # Row 9: 누적 순이익
    put(ws, "A9", "누적 순이익 (만원)", BODY_FONT, None, BORDER, left)
    put(ws, "B9", "=B8", BODY_FONT, ACCENT_FILL, BORDER, center)
    for m in range(2, 13):
        col = get_column_letter(1 + m)
        prev = get_column_letter(m)
        put(ws, f"{col}9", f"={prev}9+{col}8", BODY_FONT, ACCENT_FILL, BORDER, center)
    put(ws, "N9", "", None, None, BORDER)

    # Row 10: ROI %  = 누적 순이익 / 누적 비용
    put(ws, "A10", "ROI %", BODY_FONT, None, BORDER, left)
    put(ws, "B10", '=IF(SUM($B7:B7)=0,"-",SUM($B8:B8)/SUM($B7:B7))',
        BODY_FONT, ACCENT_FILL, BORDER, center)
    for m in range(2, 13):
        col = get_column_letter(1 + m)
        put(ws, f"{col}10",
            f'=IF(SUM($B7:{col}7)=0,"-",SUM($B8:{col}8)/SUM($B7:{col}7))',
            BODY_FONT, ACCENT_FILL, BORDER, center)
    # format as %
    for m in range(1, 13):
        col = get_column_letter(1 + m)
        ws[f"{col}10"].number_format = "0.0%"

    put(ws, "N10", "", None, None, BORDER)

    note_row = 13
    put(ws, f"A{note_row}",
        "월별 측정값을 채워 넣으세요. 회수는 누적 순이익(9행)이 0을 넘는 월을 회수 시점으로 봅니다.",
        HINT_FONT)
    ws.merge_cells(f"A{note_row}:N{note_row}")

    put(ws, f"A{note_row+1}",
        "노랑 셀은 자동 계산됩니다. 절감액·비용 두 행만 직접 입력하세요.",
        HINT_FONT)
    ws.merge_cells(f"A{note_row+1}:N{note_row+1}")

    footer_row = note_row + 3
    put(ws, f"A{footer_row}", FOOTER, FOOTER_FONT)
    ws.merge_cells(f"A{footer_row}:N{footer_row}")

    out = os.path.join(OUTDIR, "payback-model.xlsx")
    wb.save(out)
    print(f"  wrote {out}")


# ---------- 4) AI usage policy (DOCX) ----------

def _add_page_number_footer(doc, footer_text):
    """Add page numbering + custom footer text to the document."""
    section = doc.sections[0]
    footer = section.footer
    # Clear default paragraph
    p = footer.paragraphs[0]
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Add custom text run
    run = p.add_run(f"{footer_text}    ·    ")
    run.font.size = Pt(9)
    run.font.italic = True
    run.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

    # Add PAGE field
    fld_run = p.add_run()
    fld_run.font.size = Pt(9)
    fld_run.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

    fldChar1 = OxmlElement("w:fldChar")
    fldChar1.set(qn("w:fldCharType"), "begin")
    instrText = OxmlElement("w:instrText")
    instrText.set(qn("xml:space"), "preserve")
    instrText.text = "PAGE"
    fldChar2 = OxmlElement("w:fldChar")
    fldChar2.set(qn("w:fldCharType"), "end")
    fld_run._r.append(fldChar1)
    fld_run._r.append(instrText)
    fld_run._r.append(fldChar2)

    fld_run.add_text(" / ")

    fldChar3 = OxmlElement("w:fldChar")
    fldChar3.set(qn("w:fldCharType"), "begin")
    instrText2 = OxmlElement("w:instrText")
    instrText2.set(qn("xml:space"), "preserve")
    instrText2.text = "NUMPAGES"
    fldChar4 = OxmlElement("w:fldChar")
    fldChar4.set(qn("w:fldCharType"), "end")
    fld_run._r.append(fldChar3)
    fld_run._r.append(instrText2)
    fld_run._r.append(fldChar4)


def make_usage_policy():
    doc = Document()

    # Page margins
    for section in doc.sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1)
        section.right_margin = Inches(1)

    # Title
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("사내 AI 사용 정책 v0.1 (템플릿)")
    run.bold = True
    run.font.size = Pt(22)
    run.font.color.rgb = RGBColor(0x11, 0x18, 0x27)

    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    mr = meta.add_run("문서번호: [POL-AI-001]    ·    적용일: [YYYY-MM-DD]    ·    소관부서: [정보보호팀]")
    mr.font.size = Pt(10)
    mr.font.italic = True
    mr.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

    doc.add_paragraph()

    # 1. 목적과 범위
    doc.add_heading("1. 목적과 범위", level=1)
    doc.add_paragraph(
        "[회사명]은 AI 도구의 안전하고 효과적인 사용을 위해 본 정책을 수립합니다. "
        "본 정책은 사내·외부에서 사용되는 모든 생성형 AI 도구 및 LLM 기반 서비스에 적용되며, "
        "기밀 데이터 보호, 법적 리스크 최소화, 산출물 품질 확보를 목적으로 합니다."
    )

    # 2. 적용 대상
    doc.add_heading("2. 적용 대상", level=1)
    doc.add_paragraph(
        "전 임직원 및 [범위 — 정규직 / 계약직 / 외부 파트너 / 외주 인력 / 인턴] 중 "
        "사내 시스템 또는 사내 데이터에 접근하는 모든 인력에 적용됩니다."
    )

    # 3. 데이터 등급 분류
    doc.add_heading("3. 데이터 등급 분류", level=1)
    doc.add_paragraph(
        "회사가 보유한 모든 데이터는 4개 등급으로 분류됩니다. 각 등급의 정의는 사내 정보보호 규정에 따릅니다."
    )

    tbl = doc.add_table(rows=5, cols=3)
    tbl.style = "Light Grid Accent 1"
    hdr_cells = tbl.rows[0].cells
    hdr_cells[0].text = "등급"
    hdr_cells[1].text = "명칭"
    hdr_cells[2].text = "정의 (사내 채워 넣기)"
    for cell in hdr_cells:
        for p in cell.paragraphs:
            for r in p.runs:
                r.bold = True

    grades = [
        ("Lv1", "공개", "[ 예: 보도자료, 공식 웹사이트 콘텐츠 등 외부 공개 가능 정보 ]"),
        ("Lv2", "내부", "[ 예: 사내 공지, 일반 업무 자료 등 사내 한정 정보 ]"),
        ("Lv3", "기밀", "[ 예: 고객 데이터, 미공개 재무, 전략 문서 등 부서/프로젝트 한정 ]"),
        ("Lv4", "극비", "[ 예: M&A, 미공개 인사, 법적 보호 대상 등 임원 한정 ]"),
    ]
    for i, (lv, name, defn) in enumerate(grades, start=1):
        row = tbl.rows[i].cells
        row[0].text = lv
        row[1].text = name
        row[2].text = defn

    # 4. 승인된 도구
    doc.add_heading("4. 승인된 도구", level=1)
    doc.add_paragraph(
        "아래 도구만 사내 업무에 사용할 수 있습니다. 신규 도구 도입은 정보보호팀의 승인을 거쳐야 하며, "
        "각 도구별 사용 가능 데이터 등급을 준수해야 합니다."
    )

    tools_tbl = doc.add_table(rows=6, cols=4)
    tools_tbl.style = "Light Grid Accent 1"
    hdrs = ["도구", "유형", "사용 가능 등급", "비고"]
    for i, h in enumerate(hdrs):
        c = tools_tbl.rows[0].cells[i]
        c.text = h
        for p in c.paragraphs:
            for r in p.runs:
                r.bold = True

    tools = [
        ("Claude", "외부 LLM", "Lv1, Lv2 (가공 후)", "Anthropic 제공"),
        ("ChatGPT", "외부 LLM", "Lv1, Lv2 (가공 후)", "OpenAI 제공 · Enterprise 권장"),
        ("Gemini (Google Workspace)", "외부 LLM (테넌트)", "Lv1, Lv2, Lv3 (익명화)", "사내 도메인 연동 시"),
        ("Microsoft 365 Copilot", "외부 LLM (테넌트)", "Lv1, Lv2, Lv3 (익명화)", "사내 테넌트 연동 시"),
        ("[사내 LLM 명칭]", "사내 호스팅", "Lv1 ~ Lv4 (Lv4는 승인 필요)", "Air-gapped 환경"),
    ]
    for i, (name, kind, lvl, memo) in enumerate(tools, start=1):
        row = tools_tbl.rows[i].cells
        row[0].text = name
        row[1].text = kind
        row[2].text = lvl
        row[3].text = memo

    # 5. 금지 행위
    doc.add_heading("5. 금지 행위", level=1)
    doc.add_paragraph(
        "다음 행위는 본 정책에 의해 명시적으로 금지되며, 위반 시 사내 징계 절차에 따릅니다."
    )
    prohibitions = [
        "극비(Lv4) 데이터를 외부 LLM(Claude/ChatGPT/Gemini 등)에 입력하는 행위",
        "AI 출력물을 사실 검증 없이 외부(고객·언론·규제기관)에 송출하는 행위",
        "고객 개인정보(이름·연락처·주민번호·금융정보 등)를 익명화 없이 AI에 입력하는 행위",
        "회사 소스코드, 미공개 알고리즘, 영업비밀을 외부 LLM에 입력하는 행위",
        "AI 도구의 출력물을 본인이 작성한 것처럼 외부에 표시하는 행위(이해상충 발생 시)",
    ]
    for item in prohibitions:
        p = doc.add_paragraph(item, style="List Bullet")

    # 6. 외부 송출 검토 절차
    doc.add_heading("6. 외부 송출 검토 절차", level=1)
    doc.add_paragraph(
        "AI가 생성한 콘텐츠를 외부에 송출하기 전, 다음 RGSB 4관점 또는 사내 검토 체크리스트를 적용합니다."
    )
    rgsb = [
        ("R (Reason · 이성)", "논리·근거·결론의 정합성을 점검한다."),
        ("D (Domain · 분야 전문)", "해당 분야의 사실 관계와 디테일을 검증한다."),
        ("E (Empathy · 공감)", "수신자의 감정과 맥락이 깨지지 않는지 확인한다."),
        ("R (Reader · 독자)", "실제 독자의 첫 행동이 가능한지 점검한다."),
    ]
    for label, desc in rgsb:
        p = doc.add_paragraph(style="List Bullet")
        r1 = p.add_run(f"{label} — ")
        r1.bold = True
        p.add_run(desc)

    # 7. 사고 발생 시 대응
    doc.add_heading("7. 사고 발생 시 대응", level=1)
    doc.add_paragraph(
        "AI 사용 중 데이터 유출, 잘못된 외부 송출, 법적 리스크가 발생한 경우 아래 절차를 따릅니다."
    )
    incident = [
        "즉시 보고 — 발견 즉시 직속 상급자 및 정보보호팀에 보고 (24시간 이내)",
        "영향 평가 — 유출 데이터의 등급, 범위, 수신자를 평가 (72시간 이내)",
        "외부 통보 — 개인정보 유출 시 관련 법령에 따라 정보주체 및 규제기관에 통보",
        "재발 방지 — 원인 분석 및 정책·교육·도구 설정 보완",
    ]
    for i, item in enumerate(incident, start=1):
        p = doc.add_paragraph(f"{i}) {item}", style="List Number" if False else None)
        doc.paragraphs[-1].paragraph_format.left_indent = Inches(0.25)

    # 8. 교육과 인증
    doc.add_heading("8. 교육과 인증", level=1)
    doc.add_paragraph(
        "신규 입사자는 입사 후 [14일] 이내 본 정책에 대한 의무 교육을 이수해야 합니다. "
        "전 임직원은 분기 1회 정책 업데이트 및 사례 공유 교육에 참여해야 하며, "
        "교육 이수 여부는 인사 평가에 반영됩니다."
    )

    # 9. 개정 이력
    doc.add_heading("9. 개정 이력", level=1)
    hist_tbl = doc.add_table(rows=2, cols=4)
    hist_tbl.style = "Light Grid Accent 1"
    hh = hist_tbl.rows[0].cells
    hh[0].text = "버전"
    hh[1].text = "개정일"
    hh[2].text = "작성자"
    hh[3].text = "변경 사항"
    for c in hh:
        for p in c.paragraphs:
            for r in p.runs:
                r.bold = True
    row1 = hist_tbl.rows[1].cells
    row1[0].text = "v0.1"
    row1[1].text = "2026-MM-DD"
    row1[2].text = "[작성자 이름]"
    row1[3].text = "최초 작성 (템플릿 기반)"

    # Footer
    _add_page_number_footer(doc, FOOTER)

    out = os.path.join(OUTDIR, "ai-usage-policy.docx")
    doc.save(out)
    print(f"  wrote {out}")


# ---------- 5) Data classification matrix ----------

def make_data_classification():
    wb = Workbook()
    ws = wb.active
    ws.title = "데이터 등급별 사용 가능 도구"

    set_col_widths(ws, [14, 30, 14, 14, 18, 18, 14])

    put(ws, "A1", "데이터 등급별 AI 도구 사용 매트릭스",
        Font(name="Inter", size=16, bold=True, color="111827"))
    ws.merge_cells("A1:G1")
    ws.row_dimensions[1].height = 26

    put(ws, "A3", "회사명", SUBHEADER_FONT, SUBHEADER_FILL, BORDER)
    put(ws, "B3", "", BODY_FONT, None, BORDER)
    ws.merge_cells("B3:D3")
    put(ws, "E3", "작성일", SUBHEADER_FONT, SUBHEADER_FILL, BORDER)
    put(ws, "F3", "", BODY_FONT, None, BORDER)
    ws.merge_cells("F3:G3")

    # Header
    headers = ["데이터 등급", "정의", "Claude (외부)", "ChatGPT (외부)",
               "Gemini Workspace", "M365 Copilot", "사내 LLM"]
    for i, h in enumerate(headers):
        col = get_column_letter(i + 1)
        put(ws, f"{col}5", h, HEADER_FONT, HEADER_FILL, BORDER,
            Alignment(horizontal="center", vertical="center", wrap_text=True))
    ws.row_dimensions[5].height = 36

    OK = "OK"
    COND = "조건부"
    NO = "금지"
    OK_FILL = PatternFill("solid", fgColor="D1FAE5")
    COND_FILL = PatternFill("solid", fgColor="FEF3C7")
    NO_FILL = PatternFill("solid", fgColor="FECACA")

    rows = [
        ("공개 (Lv1)", "외부 공개 가능",
         (OK, ""), (OK, ""), (OK, ""), (OK, ""), (OK, "")),
        ("내부 (Lv2)", "사내 한정",
         (COND, "가공 후"), (COND, "가공 후"), (OK, ""), (OK, ""), (OK, "")),
        ("기밀 (Lv3)", "부서/프로젝트 한정",
         (NO, ""), (NO, ""), (COND, "익명화"), (COND, "익명화"), (OK, "")),
        ("극비 (Lv4)", "임원 한정 / 법적 보호",
         (NO, ""), (NO, ""), (NO, ""), (NO, ""), (COND, "승인 필요")),
    ]

    center = Alignment(horizontal="center", vertical="center", wrap_text=True)
    left = Alignment(horizontal="left", vertical="center", wrap_text=True)

    def cell_for(val):
        status, note = val
        if status == OK:
            return f"OK", OK_FILL
        elif status == COND:
            return f"조건부\n({note})" if note else "조건부", COND_FILL
        else:
            return "금지", NO_FILL

    for i, row in enumerate(rows):
        r = 6 + i
        grade, defn = row[0], row[1]
        put(ws, f"A{r}", grade, SUBHEADER_FONT, SUBHEADER_FILL, BORDER, center)
        put(ws, f"B{r}", defn, BODY_FONT, None, BORDER, left)
        for j, val in enumerate(row[2:], start=3):
            col = get_column_letter(j)
            txt, fill = cell_for(val)
            put(ws, f"{col}{r}", txt, BODY_FONT, fill, BORDER, center)
        ws.row_dimensions[r].height = 36

    # Legend
    legend_row = 6 + len(rows) + 2
    put(ws, f"A{legend_row}", "범례", SUBHEADER_FONT, SUBHEADER_FILL, BORDER, center)
    put(ws, f"B{legend_row}", "OK = 사용 가능   ·   조건부 = 가공/익명화/승인 필요   ·   금지 = 사용 불가",
        BODY_FONT, None, BORDER, left)
    ws.merge_cells(f"B{legend_row}:G{legend_row}")

    note_row = legend_row + 2
    put(ws, f"A{note_row}",
        "이 표는 표준 권장안입니다. 사내 데이터 카탈로그를 먼저 작성한 뒤 등급을 매기세요. "
        "도구별 보안 인증(SOC2, ISO27001, KISA-ISMS 등) 및 사내 테넌트 계약 조건에 따라 등급은 달라질 수 있습니다.",
        HINT_FONT, align=Alignment(wrap_text=True, vertical="top"))
    ws.merge_cells(f"A{note_row}:G{note_row}")
    ws.row_dimensions[note_row].height = 48

    footer_row = note_row + 2
    put(ws, f"A{footer_row}", FOOTER, FOOTER_FONT)
    ws.merge_cells(f"A{footer_row}:G{footer_row}")

    out = os.path.join(OUTDIR, "data-classification.xlsx")
    wb.save(out)
    print(f"  wrote {out}")


# ---------- main ----------

if __name__ == "__main__":
    os.makedirs(OUTDIR, exist_ok=True)
    make_roi_calculator()
    make_impact_matrix()
    make_payback_model()
    make_usage_policy()
    make_data_classification()
    print("Generated 5 P2 assets")
